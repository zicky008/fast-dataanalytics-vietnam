"""
DataAnalytics Vietnam - AI-Powered Data Cleaning & Dashboard Builder
Ứng dụng phân tích dữ liệu tự động với Gemini AI
"""

import streamlit as st
import pandas as pd
import numpy as np
import plotly.express as px
import plotly.graph_objects as go
from plotly.subplots import make_subplots
import io
import base64
from datetime import datetime
from typing import Dict, List, Tuple, Any

# Import utility modules
import sys
import os
sys.path.insert(0, os.path.join(os.path.dirname(__file__), '..'))
from utils.validators import safe_file_upload, sanitize_column_names, validate_dataframe
from utils.error_handlers import rate_limit_handler, user_friendly_error
from utils.performance import log_performance, PerformanceMonitor
from google import genai
from google.genai import types
import os

# IMPORTANT: KEEP THIS COMMENT
# This uses the blueprint:python_gemini integration
# Using the new google-genai SDK (renamed from google-generativeai)
# Gemini 2.5-flash and gemini-2.5-pro are the newest models

# Cấu hình trang
st.set_page_config(
    page_title="DataAnalytics Vietnam",
    page_icon="📊",
    layout="wide",
    initial_sidebar_state="expanded"
)

# CSS tùy chỉnh với gradient xanh dương-tím
st.markdown("""
<style>
    .main-header {
        background: linear-gradient(135deg, #667eea 0%, #764ba2 100%);
        padding: 2rem;
        border-radius: 10px;
        text-align: center;
        color: white;
        margin-bottom: 2rem;
    }
    .main-header h1 {
        font-size: 2.5rem;
        font-weight: 700;
        margin-bottom: 0.5rem;
    }
    .main-header p {
        font-size: 1.1rem;
        opacity: 0.9;
    }
    .step-card {
        background: white;
        padding: 1.5rem;
        border-radius: 10px;
        box-shadow: 0 2px 8px rgba(0,0,0,0.1);
        margin-bottom: 1.5rem;
        border-left: 4px solid #667eea;
    }
    .metric-card {
        background: linear-gradient(135deg, #667eea 0%, #764ba2 100%);
        padding: 1.5rem;
        border-radius: 10px;
        color: white;
        text-align: center;
    }
    .success-box {
        background: #10B981;
        color: white;
        padding: 1rem;
        border-radius: 8px;
        margin: 1rem 0;
    }
    .warning-box {
        background: #F59E0B;
        color: white;
        padding: 1rem;
        border-radius: 8px;
        margin: 1rem 0;
    }
    .info-box {
        background: #3B82F6;
        color: white;
        padding: 1rem;
        border-radius: 8px;
        margin: 1rem 0;
    }
    .stTabs [data-baseweb="tab-list"] {
        gap: 2rem;
    }
    .stTabs [data-baseweb="tab"] {
        height: 50px;
        padding-left: 20px;
        padding-right: 20px;
        background-color: #f0f2f6;
        border-radius: 5px 5px 0 0;
    }
    .stTabs [aria-selected="true"] {
        background: linear-gradient(135deg, #667eea 0%, #764ba2 100%);
        color: white;
    }
</style>
""", unsafe_allow_html=True)

# Header
st.markdown("""
<div class="main-header">
    <h1>📊 DataAnalytics Vietnam</h1>
    <p>Phân tích dữ liệu tự động với AI - Data Cleaning & Dashboard Builder</p>
</div>
""", unsafe_allow_html=True)

# Initialize session state - OPTIMIZED (single dict approach)
DEFAULT_SESSION_STATE = {
    'data': None,
    'cleaned_data': None,
    'date_cols': [],
    'categorical_cols': [],
    'numerical_cols': [],
    'dataset_description': "",
    'analysis_goal': "",
    'dashboard_blueprint': None,
    'blueprint_structured_data': None,
    'eda_results': {},
    'cleaning_log': [],
    'executive_summary': None,
    'dashboard_insights': None,
    'regenerate_insights_flag': False,
    'per_chart_insights': {},
    'regenerate_chart_insights_flag': False,
    'active_filters': {},
    'dashboard_template': None
}

# Initialize all session state keys at once (more efficient)
for key, default_value in DEFAULT_SESSION_STATE.items():
    if key not in st.session_state:
        st.session_state[key] = default_value

# Khởi tạo Gemini AI Client
def init_gemini():
    """Khởi tạo Gemini AI"""
    api_key = os.getenv('GEMINI_API_KEY')
    if api_key:
        try:
            client = genai.Client(api_key=api_key)
            return client
        except Exception as e:
            st.error(f"❌ Lỗi khởi tạo Gemini: {str(e)}")
            return None
    else:
        st.warning("⚠️ GEMINI_API_KEY chưa được cấu hình. Một số tính năng AI sẽ không khả dụng.")
        return None

# Generate AI insights với Gemini
@rate_limit_handler(max_retries=3, backoff_base=2)
@log_performance("Gemini AI Insight Generation")
def generate_ai_insight(client, prompt: str, temperature: float = 0.7, max_tokens: int = 4096) -> Tuple[bool, str]:
    """
    Tạo insights từ Gemini AI với rate limit handling và performance logging
    Returns: (success: bool, result: str)
    """
    try:
        response = client.models.generate_content(
            model='gemini-2.5-flash',
            contents=prompt,
            config=types.GenerateContentConfig(
                temperature=temperature,
                max_output_tokens=max_tokens,
            )
        )
        
        # Debug: Check response object
        if hasattr(response, 'text') and response.text:
            return (True, response.text)
        elif hasattr(response, 'candidates') and response.candidates:
            # Try to extract text from candidates
            if hasattr(response.candidates[0], 'content'):
                parts = response.candidates[0].content.parts
                text = ''.join([part.text for part in parts if hasattr(part, 'text')])
                if text:
                    return (True, text)
        
        # If we reach here, response has no usable content
        return (False, f"AI trả về phản hồi trống. Response object: {type(response)}. Có thể do prompt quá dài hoặc content bị filter.")
    except Exception as e:
        # Use user_friendly_error() for better error messages
        return (False, user_friendly_error(e))

def parse_blueprint_structured_data(metrics_text: str, layout_text: str) -> Dict[str, Any]:
    """
    Parse OQMLB blueprint text to extract structured data for dashboard generation
    Returns: {
        'primary_kpis': [{name, priority, threshold, action_recommendations}],
        'hero_kpi': {...},  # P1 KPI
        'layout_pages': [{name, purpose, visualizations, filters}]
    }
    """
    import re
    
    result = {
        'primary_kpis': [],
        'hero_kpi': None,
        'layout_pages': []
    }
    
    # Parse PRIMARY KPIs from Metrics section
    if metrics_text:
        # Look for PRIMARY KPIs section
        primary_section_match = re.search(r'\*\*⭐ PRIMARY KPIs\*\*.*?(?=\n\*\*|$)', metrics_text, re.DOTALL | re.IGNORECASE)
        if primary_section_match:
            primary_text = primary_section_match.group(0)
            
            # Extract individual KPIs (looking for patterns like "**KPI Name**: ..." or numbered lists)
            kpi_blocks = re.split(r'\n\d+\.\s+\*\*|^\d+\.\s+\*\*', primary_text, flags=re.MULTILINE)
            
            for block in kpi_blocks[1:]:  # Skip first empty split
                kpi_data = {}
                
                # Extract KPI name (first bold text or after number)
                name_match = re.search(r'^\*\*(.+?)\*\*|^(.+?):\s', block)
                if name_match:
                    kpi_data['name'] = (name_match.group(1) or name_match.group(2)).strip()
                
                # Extract priority
                priority_match = re.search(r'Priority.*?[:：]\s*(\[?P[123]\]?)', block, re.IGNORECASE)
                if priority_match:
                    kpi_data['priority'] = priority_match.group(1).strip('[]')
                else:
                    kpi_data['priority'] = 'P2'  # Default
                
                # Extract threshold/target
                threshold_match = re.search(r'(?:Target|Threshold|Ngưỡng).*?[:：]\s*(.+?)(?:\n|$)', block, re.IGNORECASE)
                if threshold_match:
                    kpi_data['threshold'] = threshold_match.group(1).strip()
                
                # Extract action recommendations
                action_match = re.search(r'(?:Action.*?Recommendation|Nên làm).*?[:：]\s*(.+?)(?=\n\*\*|$)', block, re.DOTALL | re.IGNORECASE)
                if action_match:
                    kpi_data['action_recommendations'] = action_match.group(1).strip()
                
                if kpi_data.get('name'):
                    result['primary_kpis'].append(kpi_data)
                    
                    # Set hero KPI (P1)
                    if kpi_data.get('priority') == 'P1' and not result['hero_kpi']:
                        result['hero_kpi'] = kpi_data
    
    # Parse Layout Pages
    if layout_text:
        # Find page blocks
        page_blocks = re.findall(r'\*\*Page \[?\d+\]?:?\s*(.+?)\*\*\s*(?:-\s*\*\*Purpose\*\*:?\s*(.+?))?(?=\n\*\*Page|\n##|$)', layout_text, re.DOTALL | re.IGNORECASE)
        
        for page_match in page_blocks:
            page_name = page_match[0].strip() if page_match[0] else ''
            page_purpose = page_match[1].strip() if len(page_match) > 1 and page_match[1] else ''
            
            result['layout_pages'].append({
                'name': page_name,
                'purpose': page_purpose
            })
    
    return result

def generate_dashboard_insights(client, df: pd.DataFrame, numerical_cols: list, categorical_cols: list) -> Tuple[bool, str]:
    """
    Generate AI-powered Dashboard Insights (like Bricks.ai "Overall Insights")
    Analyzes entire dashboard and returns bullet-point insights with specific numbers
    Returns: (success: bool, insights_markdown: str)
    """
    if df is None or df.empty:
        return (False, "No data available")
    
    # Calculate key statistics for numerical columns
    num_insights = []
    if numerical_cols:
        for col in numerical_cols[:5]:  # Top 5 numerical metrics
            col_mean = df[col].mean()
            col_median = df[col].median()
            col_max = df[col].max()
            col_min = df[col].min()
            col_std = df[col].std()
            col_sum = df[col].sum()
            
            num_insights.append({
                'name': col,
                'mean': col_mean,
                'median': col_median,
                'max': col_max,
                'min': col_min,
                'std': col_std,
                'sum': col_sum
            })
    
    # Calculate categorical distributions
    cat_insights = []
    if categorical_cols:
        for col in categorical_cols[:5]:  # Top 5 categorical dimensions
            value_counts = df[col].value_counts()
            top_category = value_counts.index[0] if len(value_counts) > 0 else None
            top_count = value_counts.iloc[0] if len(value_counts) > 0 else 0
            total_categories = len(value_counts)
            
            cat_insights.append({
                'name': col,
                'top_category': top_category,
                'top_count': top_count,
                'total_categories': total_categories,
                'distribution': value_counts.head(5).to_dict()
            })
    
    # Build analysis context
    num_summary = "\n".join([
        f"- **{item['name']}**: Mean={item['mean']:,.2f}, Median={item['median']:,.2f}, Max={item['max']:,.2f}, Min={item['min']:,.2f}, Sum={item['sum']:,.0f}"
        for item in num_insights
    ])
    
    cat_summary = "\n".join([
        f"- **{item['name']}**: Top={item['top_category']} ({item['top_count']} records), Total categories={item['total_categories']}"
        for item in cat_insights
    ])
    
    # Correlations (if multiple numerical columns)
    correlation_text = ""
    if len(numerical_cols) >= 2:
        corr_matrix = df[numerical_cols].corr()
        # Find top 3 correlations
        top_corrs = []
        for i in range(len(corr_matrix.columns)):
            for j in range(i+1, len(corr_matrix.columns)):
                col1 = corr_matrix.columns[i]
                col2 = corr_matrix.columns[j]
                corr_val = corr_matrix.iloc[i, j]
                if abs(corr_val) > 0.3:  # Only significant correlations
                    top_corrs.append((col1, col2, corr_val))
        
        top_corrs.sort(key=lambda x: abs(x[2]), reverse=True)
        if top_corrs:
            correlation_text = "**KEY CORRELATIONS:**\n" + "\n".join([
                f"- {c[0]} vs {c[1]}: {c[2]:.2f} correlation"
                for c in top_corrs[:3]
            ])
    
    prompt = f"""
BẠN LÀ DATA ANALYST EXPERT. Phân tích dashboard data và tạo "Overall Insights" giống Bricks.ai format.

**DATASET OVERVIEW:**
- Total Records: {len(df):,}
- Numerical Metrics: {len(numerical_cols)}
- Categorical Dimensions: {len(categorical_cols)}

**NUMERICAL METRICS SUMMARY:**
{num_summary}

**CATEGORICAL DISTRIBUTIONS:**
{cat_summary}

{correlation_text}

---

**YÊU CẦU OUTPUT (Vietnamese):**

Tạo "💡 AI Dashboard Insights" theo format SAU ĐÚNG:

### 💡 AI Dashboard Insights

**📊 Key Performance Indicators:**
- [Metric name] đạt [specific number with unit], [comparison context, e.g., "cao hơn trung bình 23%"]
- [Metric name] có giá trị trung bình [number], [insight về distribution/variability]
- [Top performer insight với specific numbers]

**📈 Trends & Patterns:**
- [Categorical dimension] có sự phân bố: [top category] chiếm [X%], [insight]
- [Correlation insight nếu có, e.g., "ROI tăng 1.5x khi Spend tăng gấp đôi"]
- [Outlier/anomaly nếu phát hiện được]

**🎯 Actionable Recommendations:**
- **[Priority #1]**: [Action với specific numbers, e.g., "Tăng budget cho segment A từ 10M → 15M (+50%)"]
- **[Priority #2]**: [Action]
- **[Priority #3]**: [Action]

---

**CRITICAL RULES:**
✅ EVERY bullet point MUST have SPECIFIC NUMBERS (no vague statements like "performance is good")
✅ Use Vietnamese business language (ROI, conversion rate, clicks, impressions, etc.)
✅ Format like Bricks.ai: bullet points, concise, data-driven
✅ Focus on TOP 3-5 insights per section (avoid information overload)
✅ Include comparative context: percentages, comparisons, benchmarks
✅ Be ACTIONABLE: recommendations must be specific and implementable
✅ Limit to 8-12 bullet points total (keep it scannable)

**EXAMPLE GOOD OUTPUT (reference format):**

### 💡 AI Dashboard Insights

**📊 Key Performance Indicators:**
- Total Spend đạt 7.86M, với Average ROI là 172.99M (gấp 22x chi phí)
- Total Clicks đạt 18M lượt, tạo ra 57M impressions (tỷ lệ engagement 31.5%)
- Conversion rate trung bình 0.085, với độ lệch chuẩn thấp (0.004) cho thấy consistency tốt

**📈 Trends & Patterns:**
- Fashion segment chiếm 50.8% total clicks (9.35M), cao hơn Food segment 3.4%
- Women 25-34 có conversion rate cao nhất (0.09), outperform nam giới 12.5%
- Pinterest channel có ROI 704.57M, gấp 175x các channel khác (Facebook: 4.03M, Twitter: 4.06M)

**🎯 Actionable Recommendations:**
- **Tăng budget Pinterest ngay**: Channel này có ROI cao nhất nhưng spend thấp nhất (1.98M), tăng lên 4M (+100%) để scale winning campaigns
- **Optimize targeting Women 25-34**: Segment này có conversion 0.09 vs 0.08 nam giới, nên shift 20-30% budget sang đây
- **Review Fashion campaigns**: Segment tạo nhiều clicks nhất, cần A/B test thêm creatives để push conversion rate từ 0.085 → 0.10 (+17%)

Generate insights now (Vietnamese):
"""
    
    return generate_ai_insight(client, prompt, temperature=0.4, max_tokens=4096)

def generate_per_chart_insights(client, df: pd.DataFrame, numerical_cols: list, categorical_cols: list) -> Tuple[bool, dict]:
    """
    Generate AI-powered insights for EACH chart/visual (Bricks.ai style)
    Uses 1 batched Gemini call to analyze ALL charts at once (lean cost)
    Returns: (success: bool, insights_dict: {chart_id: insights_text})
    """
    if df is None or df.empty:
        return (False, {})
    
    # Build chart definitions with data context
    chart_definitions = []
    
    # Chart 1: Trend Line Chart (first numerical column)
    if numerical_cols and len(numerical_cols) >= 1:
        col = numerical_cols[0]
        chart_definitions.append({
            'id': 'trend_chart',
            'type': 'Line Chart - Trend Analysis',
            'column': col,
            'stats': {
                'mean': df[col].mean(),
                'median': df[col].median(),
                'min': df[col].min(),
                'max': df[col].max(),
                'std': df[col].std()
            }
        })
    
    # Chart 2: Distribution Histogram (second numerical column)
    if numerical_cols and len(numerical_cols) >= 2:
        col = numerical_cols[1]
        chart_definitions.append({
            'id': 'distribution_chart',
            'type': 'Histogram - Distribution',
            'column': col,
            'stats': {
                'mean': df[col].mean(),
                'median': df[col].median(),
                'min': df[col].min(),
                'max': df[col].max(),
                'std': df[col].std()
            }
        })
    
    # Chart 3: Category Bar Chart (categorical vs numerical)
    if categorical_cols and numerical_cols:
        cat_col = categorical_cols[0]
        num_col = numerical_cols[0] if len(numerical_cols) >= 1 else numerical_cols[0]
        
        # Get top categories by numerical value
        top_cats = df.groupby(cat_col)[num_col].mean().sort_values(ascending=False).head(5)
        
        chart_definitions.append({
            'id': 'category_bar_chart',
            'type': 'Bar Chart - Category Comparison',
            'categorical_col': cat_col,
            'numerical_col': num_col,
            'top_categories': {
                str(k): float(v) for k, v in top_cats.items()
            }
        })
    
    # Chart 4: Correlation Heatmap (if 3+ numerical columns)
    if len(numerical_cols) >= 3:
        corr_matrix = df[numerical_cols[:5]].corr()
        
        # Find top correlations
        top_corrs = []
        for i in range(len(corr_matrix.columns)):
            for j in range(i+1, len(corr_matrix.columns)):
                col1 = corr_matrix.columns[i]
                col2 = corr_matrix.columns[j]
                corr_val = corr_matrix.iloc[i, j]
                if abs(corr_val) > 0.3:
                    top_corrs.append((col1, col2, corr_val))
        
        top_corrs.sort(key=lambda x: abs(x[2]), reverse=True)
        
        chart_definitions.append({
            'id': 'correlation_heatmap',
            'type': 'Heatmap - Correlation Matrix',
            'columns': numerical_cols[:5],
            'top_correlations': [
                {'col1': c[0], 'col2': c[1], 'value': float(c[2])}
                for c in top_corrs[:3]
            ]
        })
    
    # Build structured prompt for batched analysis
    charts_context = ""
    for i, chart in enumerate(chart_definitions, 1):
        charts_context += f"\n**CHART {i}: {chart['type']}**\n"
        
        if chart['id'] == 'trend_chart':
            charts_context += f"- Column: {chart['column']}\n"
            charts_context += f"- Mean: {chart['stats']['mean']:,.2f}, Median: {chart['stats']['median']:,.2f}\n"
            charts_context += f"- Range: {chart['stats']['min']:,.2f} to {chart['stats']['max']:,.2f}\n"
        
        elif chart['id'] == 'distribution_chart':
            charts_context += f"- Column: {chart['column']}\n"
            charts_context += f"- Mean: {chart['stats']['mean']:,.2f}, Std: {chart['stats']['std']:,.2f}\n"
            charts_context += f"- Range: {chart['stats']['min']:,.2f} to {chart['stats']['max']:,.2f}\n"
        
        elif chart['id'] == 'category_bar_chart':
            charts_context += f"- Categorical: {chart['categorical_col']}\n"
            charts_context += f"- Numerical: {chart['numerical_col']}\n"
            charts_context += f"- Top Categories:\n"
            for cat, val in list(chart['top_categories'].items())[:3]:
                charts_context += f"  • {cat}: {val:,.2f}\n"
        
        elif chart['id'] == 'correlation_heatmap':
            charts_context += f"- Columns: {', '.join(chart['columns'])}\n"
            charts_context += f"- Top Correlations:\n"
            for corr in chart['top_correlations']:
                charts_context += f"  • {corr['col1']} ↔ {corr['col2']}: {corr['value']:.2f}\n"
        
        charts_context += "\n"
    
    prompt = f"""
BẠN LÀ DATA VISUALIZATION EXPERT. Phân tích từng chart và tạo "Key Insights" box cho MỖI VISUAL (Bricks.ai format).

**DATASET:** {len(df):,} records

{charts_context}

---

**YÊU CẦU OUTPUT:**

Tạo insights cho TỪNG CHART theo format SAU:

### CHART_1_INSIGHTS
• [Insight #1 với specific number, e.g., "Revenue trung bình 30 triệu VND, cao hơn median 15%"]
• [Insight #2 về pattern/trend, e.g., "Peak value đạt 57.3 triệu, gấp 1.9x trung bình"]
• [Insight #3 actionable, e.g., "Tập trung vào top 25% performers để tối ưu ROI"]

### CHART_2_INSIGHTS
• [Insight #1]
• [Insight #2]
• [Insight #3]

### CHART_3_INSIGHTS
• [Insight #1]
• [Insight #2]
• [Insight #3]

### CHART_4_INSIGHTS
• [Insight #1]
• [Insight #2]

---

**CRITICAL RULES:**
✅ Mỗi chart: 2-3 bullet points NGẮN GỌN (max 1 line per bullet)
✅ EVERY bullet MUST have SPECIFIC NUMBERS (no vague statements)
✅ Vietnamese business language (ROI, conversion, clicks, revenue, etc.)
✅ Format: "### CHART_X_INSIGHTS" header + bullets
✅ Actionable insights: comparisons, benchmarks, recommendations
✅ Keep scannable: đủ thông tin nhưng không quá dài

**EXAMPLE OUTPUT:**

### CHART_1_INSIGHTS
• Total Spend trend dao động 10-50 triệu VND, với peak tại record #45 (57.3 triệu)
• Spend trung bình 28.5 triệu/campaign, ổn định với std thấp (±8.2 triệu)
• Top 20% campaigns chiếm 60% total spend, cần review allocation

### CHART_2_INSIGHTS
• ROI phân bố lệch phải, majority (65%) trong range 100-200M
• Outlier campaigns có ROI 700M+ (Pinterest channel), cần scale
• Median ROI 172M cao hơn mean 15%, cho thấy consistency tốt

Generate insights cho TỪNG CHART now (Vietnamese, specific numbers):
"""
    
    # Call Gemini API
    success, response = generate_ai_insight(client, prompt, temperature=0.4, max_tokens=6144)
    
    if not success:
        return (False, {})
    
    # Parse response into dict {chart_id: insights}
    insights_dict = {}
    
    # Extract insights for each chart
    import re
    chart_blocks = re.findall(r'### CHART_(\d+)_INSIGHTS\s*\n((?:•[^\n]+\n?)+)', response, re.MULTILINE)
    
    for chart_num, insights_text in chart_blocks:
        chart_idx = int(chart_num) - 1
        if chart_idx < len(chart_definitions):
            chart_id = chart_definitions[chart_idx]['id']
            insights_dict[chart_id] = insights_text.strip()
    
    return (True, insights_dict)

def save_dashboard_template(blueprint_data: dict, analysis_goal: str, categorical_cols: list, numerical_cols: list, date_cols: list) -> str:
    """
    Save dashboard configuration as JSON template
    Returns: JSON string
    """
    import json
    from datetime import datetime
    
    template = {
        "version": "1.0",
        "created_at": datetime.now().isoformat(),
        "analysis_goal": analysis_goal,
        "column_types": {
            "categorical": categorical_cols,
            "numerical": numerical_cols,
            "date": date_cols
        },
        "blueprint": blueprint_data,
        "template_name": f"Dashboard Template - {datetime.now().strftime('%Y%m%d_%H%M%S')}"
    }
    
    return json.dumps(template, indent=2, ensure_ascii=False)

def load_dashboard_template(template_json: str) -> Tuple[bool, dict]:
    """
    Load dashboard template from JSON
    Returns: (success: bool, template_dict: dict)
    """
    import json
    
    try:
        template = json.loads(template_json)
        
        # Validate template structure
        required_keys = ["version", "analysis_goal", "column_types", "blueprint"]
        if not all(key in template for key in required_keys):
            return (False, {"error": "Invalid template format - missing required keys"})
        
        return (True, template)
    except json.JSONDecodeError as e:
        return (False, {"error": f"JSON decode error: {str(e)}"})
    except Exception as e:
        return (False, {"error": f"Template load error: {str(e)}"})

def generate_executive_summary(client, df: pd.DataFrame, blueprint: str, analysis_goal: str = "") -> Tuple[bool, str]:
    """
    Generate AI-powered Executive Summary (1-page overview)
    Returns: (success: bool, summary_markdown: str)
    """
    if df is None or df.empty:
        return (False, "No data available")
    
    # Calculate key statistics
    total_rows = len(df)
    total_cols = len(df.columns)
    
    # Get numerical column summaries
    num_cols = df.select_dtypes(include=[np.number]).columns.tolist()
    num_summary = ""
    if num_cols:
        for col in num_cols[:5]:  # Top 5 numerical columns
            col_sum = df[col].sum()
            col_mean = df[col].mean()
            col_std = df[col].std()
            num_summary += f"\n- **{col}**: Total={col_sum:,.0f}, Average={col_mean:,.2f}, Std={col_std:,.2f}"
    
    # Get categorical column distributions
    cat_cols = df.select_dtypes(include=['object', 'category']).columns.tolist()
    cat_summary = ""
    if cat_cols:
        for col in cat_cols[:3]:  # Top 3 categorical columns
            top_vals = df[col].value_counts().head(3)
            cat_summary += f"\n- **{col}**: Top values: {', '.join([f'{k} ({v})' for k, v in top_vals.items()])}"
    
    prompt = f"""
BẠN LÀ SENIOR BUSINESS ANALYST. Tạo EXECUTIVE SUMMARY (Tóm tắt Quản lý) chuyên nghiệp cho báo cáo phân tích dữ liệu.

**DATASET OVERVIEW:**
- Rows: {total_rows:,}
- Columns: {total_cols}

**NUMERICAL METRICS:**{num_summary}

**CATEGORICAL DIMENSIONS:**{cat_summary}

**ANALYSIS GOAL:**
{analysis_goal if analysis_goal else "Phân tích tổng quan và tìm insights"}

**DASHBOARD BLUEPRINT (AI-generated):**
{blueprint[:2000] if blueprint else "No blueprint available"}

---

**YÊU CẦU OUTPUT (Vietnamese):**

Tạo 1-page Executive Summary theo format sau (MUST HAVE SPECIFIC NUMBERS):

## 📊 EXECUTIVE SUMMARY

### 🎯 Overall Performance
[1-2 câu tổng quan về hiệu suất tổng thể, có số liệu cụ thể]
- Example: "Doanh thu đạt 1.2M VND (+23% so với tháng trước), vượt target 15%"

### ✅ Top 3 Wins (Điểm Mạnh)
1. **[Metric Name]**: [Specific number] ([% change] vs benchmark)
   - Impact: [Business impact với số liệu]
2. **[Metric Name]**: [Specific number] ([% change])
   - Impact: [Business impact]
3. **[Metric Name]**: [Specific number] ([% change])
   - Impact: [Business impact]

### ⚠️ Top 2 Areas of Concern (Vấn Đề Cần Chú Ý)
1. **[Issue]**: [Current metric] ([% worse than target/benchmark])
   - Root cause: [Brief explanation]
   - Urgency: [High/Medium]
2. **[Issue]**: [Current metric] ([% deviation])
   - Root cause: [Brief explanation]
   - Urgency: [High/Medium]

### 🚀 Top 3 Recommended Actions (PRIORITY ORDER)
**[#1 - URGENT]** [Action title]
- **What**: [Specific action với số liệu, e.g., "Tăng budget từ 10M → 15M VND (+50%)"]
- **Why**: [Rationale với data evidence]
- **Expected Impact**: [Predicted outcome, e.g., "Revenue tăng 16M VND trong Q4 (+25%)"]
- **Timeline**: [When to do, e.g., "Ngay trong tuần này"]

**[#2 - HIGH]** [Action title]
- **What**: [Specific action]
- **Why**: [Rationale]
- **Expected Impact**: [Predicted outcome]
- **Timeline**: [When]

**[#3 - MEDIUM]** [Action title]
- **What**: [Specific action]
- **Why**: [Rationale]
- **Expected Impact**: [Predicted outcome]
- **Timeline**: [When]

---

**CRITICAL RULES:**
✅ MUST include SPECIFIC NUMBERS in every section (no vague statements)
✅ Use Vietnamese business terminology
✅ Be concise (1 page = ~500-600 words max)
✅ Focus on ACTIONABLE insights (what decision-makers should DO)
✅ Use comparative context (vs target, vs previous period, vs benchmark)
✅ Prioritize by business impact (Revenue > Cost > Efficiency)

Generate the Executive Summary now:
"""
    
    return generate_ai_insight(client, prompt, temperature=0.4, max_tokens=8192)

def generate_pdf_report(df: pd.DataFrame, executive_summary: str, blueprint: str, analysis_goal: str = "") -> bytes:
    """
    Generate Professional PDF Report using ReportLab
    Returns: PDF bytes for download
    """
    from reportlab.lib.pagesizes import letter, A4
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.units import inch
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, PageBreak
    from reportlab.lib import colors
    from reportlab.lib.enums import TA_CENTER, TA_LEFT
    
    buffer = io.BytesIO()
    doc = SimpleDocTemplate(buffer, pagesize=letter, rightMargin=72, leftMargin=72, topMargin=72, bottomMargin=18)
    
    # Container for PDF elements
    story = []
    
    # Styles
    styles = getSampleStyleSheet()
    title_style = ParagraphStyle(
        'CustomTitle',
        parent=styles['Heading1'],
        fontSize=24,
        textColor=colors.HexColor('#667eea'),
        spaceAfter=30,
        alignment=TA_CENTER
    )
    heading_style = ParagraphStyle(
        'CustomHeading',
        parent=styles['Heading2'],
        fontSize=16,
        textColor=colors.HexColor('#667eea'),
        spaceAfter=12,
        spaceBefore=12
    )
    normal_style = styles['Normal']
    
    # Cover Page
    story.append(Spacer(1, 2*inch))
    story.append(Paragraph("📊 DataAnalytics Vietnam", title_style))
    story.append(Paragraph("Professional Data Analytics Report", styles['Heading3']))
    story.append(Spacer(1, 0.5*inch))
    story.append(Paragraph(f"<b>Date:</b> {datetime.now().strftime('%d/%m/%Y')}", normal_style))
    story.append(Paragraph(f"<b>Dataset:</b> {len(df)} rows × {len(df.columns)} columns", normal_style))
    if analysis_goal:
        story.append(Spacer(1, 0.3*inch))
        story.append(Paragraph(f"<b>Analysis Goal:</b>", normal_style))
        story.append(Paragraph(analysis_goal, normal_style))
    story.append(PageBreak())
    
    # Executive Summary
    if executive_summary:
        story.append(Paragraph("Executive Summary", title_style))
        story.append(Spacer(1, 0.2*inch))
        
        # Parse markdown to paragraphs
        for line in executive_summary.split('\n'):
            if line.strip():
                if line.startswith('##'):
                    story.append(Paragraph(line.replace('#', '').strip(), heading_style))
                elif line.startswith('###'):
                    story.append(Paragraph(line.replace('#', '').strip(), styles['Heading3']))
                else:
                    story.append(Paragraph(line, normal_style))
                story.append(Spacer(1, 6))
        
        story.append(PageBreak())
    
    # Key Metrics Table
    story.append(Paragraph("Key Performance Indicators", title_style))
    story.append(Spacer(1, 0.2*inch))
    
    num_cols = df.select_dtypes(include=[np.number]).columns.tolist()[:6]
    if num_cols:
        kpi_data = [['Metric', 'Total', 'Average', 'Std Dev', 'Min', 'Max']]
        for col in num_cols:
            kpi_data.append([
                col,
                f"{df[col].sum():,.0f}",
                f"{df[col].mean():,.2f}",
                f"{df[col].std():,.2f}",
                f"{df[col].min():,.2f}",
                f"{df[col].max():,.2f}"
            ])
        
        kpi_table = Table(kpi_data, colWidths=[2*inch, 1*inch, 1*inch, 1*inch, 0.8*inch, 0.8*inch])
        kpi_table.setStyle(TableStyle([
            ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor('#667eea')),
            ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
            ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
            ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
            ('FONTSIZE', (0, 0), (-1, 0), 12),
            ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
            ('BACKGROUND', (0, 1), (-1, -1), colors.beige),
            ('GRID', (0, 0), (-1, -1), 1, colors.black)
        ]))
        story.append(kpi_table)
        story.append(PageBreak())
    
    # Dashboard Blueprint
    if blueprint:
        story.append(Paragraph("Dashboard Blueprint (OQMLB Framework)", title_style))
        story.append(Spacer(1, 0.2*inch))
        
        # Truncate blueprint to fit PDF
        blueprint_short = blueprint[:3000] + "..." if len(blueprint) > 3000 else blueprint
        for line in blueprint_short.split('\n'):
            if line.strip():
                if line.startswith('##'):
                    story.append(Paragraph(line.replace('#', '').strip(), heading_style))
                elif line.startswith('###'):
                    story.append(Paragraph(line.replace('#', '').strip(), styles['Heading3']))
                else:
                    story.append(Paragraph(line[:200], normal_style))  # Truncate long lines
                story.append(Spacer(1, 4))
    
    # Build PDF
    doc.build(story)
    buffer.seek(0)
    return buffer.getvalue()

def generate_powerpoint_report(df: pd.DataFrame, executive_summary: str, blueprint: str, analysis_goal: str = "") -> bytes:
    """
    Generate Professional PowerPoint Presentation using python-pptx
    Returns: PPTX bytes for download
    """
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
    
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Define color scheme
    primary_color = RGBColor(102, 126, 234)  # #667eea
    accent_color = RGBColor(118, 75, 162)    # #764ba2
    
    # Slide 1: Title Slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "📊 DataAnalytics Vietnam"
    subtitle.text = f"Professional Data Analytics Report\n{datetime.now().strftime('%d/%m/%Y')}"
    
    # Style title
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.color.rgb = primary_color
    
    # Slide 2: Dataset Overview
    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)
    title = slide.shapes.title
    body = slide.placeholders[1]
    
    title.text = "📁 Dataset Overview"
    tf = body.text_frame
    tf.text = f"Dataset Size: {len(df):,} rows × {len(df.columns)} columns"
    
    if analysis_goal:
        p = tf.add_paragraph()
        p.text = f"Analysis Goal: {analysis_goal}"
        p.level = 1
    
    # Add column summary
    p = tf.add_paragraph()
    p.text = f"Numerical Columns: {len(df.select_dtypes(include=[np.number]).columns)}"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = f"Categorical Columns: {len(df.select_dtypes(include=['object', 'category']).columns)}"
    p.level = 1
    
    # Slide 3-5: Executive Summary (split into multiple slides)
    if executive_summary:
        # Parse executive summary sections
        sections = executive_summary.split('###')
        
        for i, section in enumerate(sections[1:4]):  # Take first 3 sections (Overall, Wins, Concerns)
            slide = prs.slides.add_slide(bullet_slide_layout)
            lines = section.strip().split('\n')
            
            # First line is section title
            title = slide.shapes.title
            title.text = lines[0].strip() if lines else "Executive Summary"
            
            # Rest is content
            body = slide.placeholders[1]
            tf = body.text_frame
            tf.clear()
            
            for line in lines[1:]:
                if line.strip():
                    p = tf.add_paragraph()
                    p.text = line.strip().lstrip('-').lstrip('*').strip()
                    p.font.size = Pt(14)
                    if line.strip().startswith(('1.', '2.', '3.', '**')):
                        p.level = 0
                    else:
                        p.level = 1
    
    # Slide 6: Key Metrics Table
    slide = prs.slides.add_slide(bullet_slide_layout)
    title = slide.shapes.title
    title.text = "📊 Key Performance Indicators"
    
    num_cols = df.select_dtypes(include=[np.number]).columns.tolist()[:5]
    if num_cols:
        # Add table
        rows = len(num_cols) + 1
        cols = 4
        left = Inches(1)
        top = Inches(2)
        width = Inches(8)
        height = Inches(4)
        
        table = slide.shapes.add_table(rows, cols, left, top, width, height).table
        
        # Header row
        table.cell(0, 0).text = 'Metric'
        table.cell(0, 1).text = 'Total'
        table.cell(0, 2).text = 'Average'
        table.cell(0, 3).text = 'Std Dev'
        
        # Data rows
        for i, col in enumerate(num_cols):
            table.cell(i+1, 0).text = col
            table.cell(i+1, 1).text = f"{df[col].sum():,.0f}"
            table.cell(i+1, 2).text = f"{df[col].mean():,.2f}"
            table.cell(i+1, 3).text = f"{df[col].std():,.2f}"
    
    # Slide 7: Recommendations
    slide = prs.slides.add_slide(bullet_slide_layout)
    title = slide.shapes.title
    title.text = "🚀 Recommended Actions"
    
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.text = "Based on data analysis:"
    
    # Extract recommendations from executive summary
    if executive_summary and "Recommended Actions" in executive_summary:
        rec_section = executive_summary.split("Recommended Actions")[-1]
        rec_lines = rec_section.split('\n')[:10]  # First 10 lines
        
        for line in rec_lines:
            if line.strip() and not line.startswith('#'):
                p = tf.add_paragraph()
                p.text = line.strip().lstrip('-').lstrip('*').strip()
                p.font.size = Pt(14)
    
    # Slide 8: Thank You
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "Thank You! 🙏"
    subtitle.text = "DataAnalytics Vietnam\nPowered by Gemini AI"
    
    # Save to bytes
    buffer = io.BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return buffer.getvalue()

# Sidebar - Chọn workflow
st.sidebar.markdown("## 🎯 Quy trình làm việc")
st.sidebar.markdown("""
**OQMLB Framework:**
- **O**: Objectives (Mục tiêu)
- **Q**: Questions (Câu hỏi)
- **M**: Metrics (Chỉ số)
- **L**: Layout (Bố cục)
- **B**: Build (Xây dựng)
""")

workflow_step = st.sidebar.radio(
    "Chọn bước làm việc:",
    ["📁 Step 1: Data Import & Overview", 
     "🧹 Step 2: Data Cleaning", 
     "📊 Step 3: EDA & Insights",
     "🎨 Step 4: Dashboard Blueprint",
     "🚀 Step 5: Build Dashboard"],
    index=0
)

# Main content based on workflow step
if workflow_step == "📁 Step 1: Data Import & Overview":
    st.markdown("""
    <div class="step-card">
        <h2>📁 Bước 1: Data Import & Overview</h2>
        <p>Import dữ liệu và phân tích tổng quan</p>
    </div>
    """, unsafe_allow_html=True)
    
    col1, col2 = st.columns([2, 1])
    
    with col1:
        st.markdown("### 📤 Upload Dataset")
        uploaded_file = st.file_uploader(
            "Chọn file CSV hoặc Excel",
            type=['csv', 'xlsx', 'xls'],
            help="Tải lên file dữ liệu của bạn"
        )
        
        if uploaded_file:
            # Use safe_file_upload() utility with comprehensive error handling
            success, df, message = safe_file_upload(
                uploaded_file,
                max_size_mb=200,
                show_progress=True
            )
            
            if success:
                st.session_state.data = df
                st.success(message)
            else:
                st.error(message)
                st.stop()
            
            try:
                
                # 🚀 PHASE 2B: SMART DATA REFRESH (Bricks.ai one-click refresh!)
                if st.session_state.blueprint_structured_data:
                    st.markdown("---")
                    
                    # Auto-detect column matching first
                    existing_cols = set(st.session_state.categorical_cols + st.session_state.numerical_cols)
                    new_cols = set(df.columns)
                    matching_cols = existing_cols.intersection(new_cols)
                    match_rate = len(matching_cols) / len(existing_cols) * 100 if existing_cols else 0
                    
                    if match_rate >= 70:
                        # HIGH MATCH RATE: Show Smart Refresh option
                        st.info("🔄 **Smart Refresh detected!** Bạn đã có dashboard blueprint. Có thể refresh data nhanh!")
                        
                        refresh_col1, refresh_col2 = st.columns([3, 2])
                        
                        with refresh_col1:
                            st.markdown("#### 🚀 Quick Dashboard Refresh")
                            st.caption("Upload data mới với cùng structure → giữ nguyên blueprint → skip Steps 2-4")
                            
                            st.success(f"✅ Column matching: {match_rate:.0f}% ({len(matching_cols)}/{len(existing_cols)} columns)")
                            st.write(f"**Matched columns:** {', '.join(sorted(matching_cols)[:10])}{'...' if len(matching_cols) > 10 else ''}")
                            
                            if st.button("⚡ Smart Refresh - Jump to Dashboard", type="primary", use_container_width=True):
                                # Preserve column types and blueprint
                                st.session_state.cleaned_data = df
                                st.success("🎉 Smart Refresh complete! Jumping to Step 5...")
                                st.info("💡 Blueprint và chart configs được giữ nguyên, chỉ data được refresh!")
                                st.rerun()
                        
                        with refresh_col2:
                            st.markdown("#### 🔁 Rebuild Option")
                            st.caption("Hoặc rebuild dashboard từ đầu")
                            if st.button("📝 Rebuild Dashboard", use_container_width=True):
                                # Clear blueprint to force full rebuild
                                st.session_state.blueprint_structured_data = None
                                st.session_state.dashboard_blueprint = None
                                st.session_state.cleaned_data = None
                                st.success("✅ Cleared old blueprint - proceed with normal flow below")
                                st.rerun()
                    else:
                        # LOW MATCH RATE: Force rebuild
                        st.warning(f"⚠️ **Column matching thấp: {match_rate:.0f}%** - Data structure khác, cần rebuild dashboard!")
                        st.write(f"**Missing columns:** {', '.join(existing_cols - new_cols)}")
                        st.write(f"**New columns:** {', '.join(new_cols - existing_cols)}")
                        
                        if st.button("🔁 Clear Blueprint & Rebuild", type="primary", use_container_width=True):
                            # Force clear blueprint and reset workflow
                            st.session_state.blueprint_structured_data = None
                            st.session_state.dashboard_blueprint = None
                            st.session_state.cleaned_data = None
                            st.session_state.per_chart_insights = {}
                            st.session_state.dashboard_insights = None
                            st.success("✅ Blueprint cleared! Proceed with Steps 2-4 below")
                            st.rerun()
                        
                        st.info("💡 Rebuild sẽ cho phép bạn làm lại toàn bộ Steps 2-4 với data structure mới")
                    
                    st.markdown("---")
                
                # Mô tả dataset
                st.markdown("### 📝 Mô tả Dataset")
                st.session_state.dataset_description = st.text_area(
                    "Mô tả dataset của bạn (loại dữ liệu, ngành nghề, v.v.):",
                    value=st.session_state.dataset_description,
                    placeholder="Ví dụ: Dữ liệu bán hàng tháng 1-12/2024, bao gồm thông tin khách hàng, sản phẩm, doanh thu..."
                )
                
                st.session_state.analysis_goal = st.text_area(
                    "Mục tiêu phân tích:",
                    value=st.session_state.analysis_goal,
                    placeholder="Ví dụ: Phân tích xu hướng doanh thu, tìm sản phẩm bán chạy, phân khúc khách hàng..."
                )
                
                # Tự động phân loại columns
                st.markdown("### 🔍 Phân loại Columns")
                st.info("📌 **Lưu ý**: Hãy xác định chính xác loại dữ liệu để AI phân tích đúng!")
                
                all_cols = df.columns.tolist()
                
                # Auto-detect datetime columns
                auto_date_cols = []
                if len(df) > 0:
                    for col in df.columns:
                        if df[col].dtype == 'datetime64[ns]':
                            auto_date_cols.append(col)
                        elif df[col].dtype == 'object':
                            try:
                                pd.to_datetime(df[col], errors='coerce')
                                non_null_ratio = df[col].notna().sum() / len(df)
                                if non_null_ratio > 0.8:
                                    sample_parsed = pd.to_datetime(df[col].dropna().head(100), errors='coerce')
                                    if len(sample_parsed) > 0 and sample_parsed.notna().sum() / len(sample_parsed) > 0.7:
                                        auto_date_cols.append(col)
                            except:
                                pass
                
                # Auto-detect categorical and numerical (excluding detected dates)
                remaining_cols = [col for col in df.columns if col not in auto_date_cols]
                auto_categorical = [col for col in df.select_dtypes(include=['object', 'category']).columns if col not in auto_date_cols]
                auto_numerical = [col for col in df.select_dtypes(include=['int64', 'float64']).columns if col not in auto_date_cols]
                
                st.session_state.date_cols = st.multiselect(
                    "📅 Date/Time Columns (ngày tháng, thời gian):",
                    options=all_cols,
                    default=auto_date_cols
                )
                
                st.session_state.categorical_cols = st.multiselect(
                    "📂 Categorical Columns (phân loại, text, ID, tên...):",
                    options=all_cols,
                    default=auto_categorical
                )
                
                st.session_state.numerical_cols = st.multiselect(
                    "📊 Numerical Columns (số học, để tính toán):",
                    options=all_cols,
                    default=auto_numerical
                )
                
                # Parse datetime columns
                if st.session_state.date_cols:
                    for col in st.session_state.date_cols:
                        try:
                            df[col] = pd.to_datetime(df[col], errors='coerce')
                        except:
                            st.warning(f"⚠️ Không thể parse column '{col}' thành datetime")
                    st.session_state.data = df
                
                # Data overview
                if st.button("🔍 Phân tích tổng quan", type="primary"):
                    with st.spinner("Đang phân tích..."):
                        st.markdown("### 📊 Data Overview")
                        
                        # Basic info
                        col_a, col_b, col_c = st.columns(3)
                        with col_a:
                            st.metric("Tổng số dòng", f"{len(df):,}")
                        with col_b:
                            st.metric("Tổng số cột", len(df.columns))
                        with col_c:
                            st.metric("Missing values", df.isnull().sum().sum())
                        
                        # Data types
                        st.markdown("#### 📋 Kiểu dữ liệu")
                        st.dataframe(pd.DataFrame({
                            'Column': df.columns,
                            'Type': df.dtypes.astype(str),
                            'Non-Null': df.count().values,
                            'Null': df.isnull().sum().values,
                            'Null %': (df.isnull().sum().values.astype(float) / len(df) * 100).round(2)
                        }), width='stretch')
                        
                        # Preview
                        st.markdown("#### 👀 Xem trước dữ liệu (5 dòng đầu)")
                        st.dataframe(df.head(), width='stretch')
                        
                        # Summary statistics
                        if st.session_state.numerical_cols:
                            st.markdown("#### 📈 Thống kê mô tả (Numerical)")
                            st.dataframe(df[st.session_state.numerical_cols].describe(), width='stretch')
                        
                        # AI Insights với ISO 8000 Framework
                        client = init_gemini()
                        if client:
                            st.markdown("#### 🤖 AI Data Quality Assessment (ISO 8000 Framework)")
                            with st.spinner("Gemini đang phân tích theo chuẩn ISO 8000..."):
                                # ISO 8000 comprehensive prompt
                                missing_pct = (df.isnull().sum().sum() / (len(df) * len(df.columns)) * 100).round(2)
                                duplicates_count = df.duplicated().sum()
                                
                                prompt = f"""
Thực hiện **DATA QUALITY ASSESSMENT** theo chuẩn **ISO 8000** cho dataset sau:

**DATASET CONTEXT:**
- Tên dataset: {st.session_state.dataset_description}
- Mục tiêu phân tích: {st.session_state.analysis_goal}
- Kích thước: {len(df):,} rows × {len(df.columns)} columns
- Date/Time columns: {st.session_state.date_cols}
- Categorical columns ({len(st.session_state.categorical_cols)}): {', '.join(st.session_state.categorical_cols[:5])}{'...' if len(st.session_state.categorical_cols) > 5 else ''}
- Numerical columns ({len(st.session_state.numerical_cols)}): {', '.join(st.session_state.numerical_cols[:5])}{'...' if len(st.session_state.numerical_cols) > 5 else ''}

**DATA QUALITY METRICS:**
- Missing values: {df.isnull().sum().sum()} cells ({missing_pct}% of total data)
- Duplicate rows: {duplicates_count} ({(duplicates_count/len(df)*100):.1f}%)
- Memory usage: {(df.memory_usage(deep=True).sum() / 1024**2):.2f} MB

**YÊU CẦU PHÂN TÍCH (ISO 8000 Framework):**

Đánh giá dataset theo **6 chiều chất lượng dữ liệu**:

**1. ACCURACY (Độ chính xác):**
- Có cột nào có giá trị không hợp lý? (tuổi âm, ngày tương lai, giá trị outlier bất thường)
- Format có nhất quán? (số điện thoại, email, currency)

**2. COMPLETENESS (Đầy đủ):**
- Missing values ở mức nào? (<5% OK, 5-20% Medium, >20% Critical)
- Cột nào thiếu nhiều data nhất? Có thể drop hay cần impute?

**3. CONSISTENCY (Nhất quán):**
- Có giá trị trùng lặp hay mâu thuẫn? (VD: "USA" vs "United States")
- Format có đồng nhất? (date formats, text casing)

**4. TIMELINESS (Kịp thời):**
- Data có còn relevant với mục tiêu phân tích?
- Time range có phù hợp?

**5. VALIDITY (Hợp lệ):**
- Giá trị có nằm trong range chấp nhận được?
- Data types có đúng? (numerical stored as text, dates as strings)

**6. UNIQUENESS (Duy nhất):**
- Có duplicate records? Cần xử lý thế nào?

**OUTPUT YÊU CẦU:**

## 📊 Tóm tắt Chất lượng (Quality Score: X/100)

✅ **Điểm mạnh** (1-3 points)
⚠️ **Vấn đề cần lưu ý** (3-5 issues chi tiết, ranked theo priority)
🔧 **Cleaning Recommendations** (Step-by-step actions cho Step 2)

## 🎯 Next Steps

**Ưu tiên cao:**
1. [Action item 1]
2. [Action item 2]

**Ưu tiên trung bình:**
- [Action item 3]

**Gợi ý thêm:**
- [Nice-to-have improvements]

---

**FORMAT:** Markdown, tiếng Việt, professional, cụ thể, actionable. Prioritize critical issues. Đưa ra số liệu cụ thể khi có thể.
"""
                                success, insight = generate_ai_insight(client, prompt, temperature=0.3, max_tokens=8192)
                                if success:
                                    with st.expander("📋 Full ISO 8000 Assessment Report", expanded=True):
                                        st.markdown(insight)
                                else:
                                    st.error(f"❌ {insight}")
                        
            except Exception as e:
                st.error(f"❌ Lỗi đọc file: {str(e)}")
    
    with col2:
        st.markdown("### 💡 Hướng dẫn")
        st.markdown("""
        **Bước 1 giúp bạn:**
        - ✅ Import dữ liệu vào hệ thống
        - ✅ Hiểu cấu trúc dataset
        - ✅ Phân loại columns chính xác
        - ✅ Phát hiện vấn đề sớm
        
        **Lưu ý:**
        - Phải mô tả rõ dataset
        - Xác định đúng categorical/numerical
        - Kiểm tra missing values
        """)
        
        if st.session_state.data is not None:
            st.markdown("### ✅ Trạng thái")
            st.success(f"Dataset đã load: {len(st.session_state.data)} rows")
            st.info(f"Categorical: {len(st.session_state.categorical_cols)} cột")
            st.info(f"Numerical: {len(st.session_state.numerical_cols)} cột")

elif workflow_step == "🧹 Step 2: Data Cleaning":
    st.markdown("""
    <div class="step-card">
        <h2>🧹 Bước 2: Data Cleaning</h2>
        <p>8 bước làm sạch dữ liệu chuyên nghiệp</p>
    </div>
    """, unsafe_allow_html=True)
    
    if st.session_state.data is None:
        st.warning("⚠️ Vui lòng hoàn thành Step 1 trước!")
    else:
        df = st.session_state.data.copy()
        
        st.markdown("### 🎯 8 Bước Data Cleaning")
        st.info("👇 Mở rộng từng bước bên dưới để thực hiện làm sạch dữ liệu")
        
        # Step 1: Relevance
        with st.expander("1️⃣ Data Relevance - Loại bỏ cột/dòng không cần thiết", expanded=False):
            cols_to_keep = st.multiselect(
                "Chọn các cột cần giữ lại:",
                options=df.columns.tolist(),
                default=df.columns.tolist()
            )
            if cols_to_keep:
                df = df[cols_to_keep]
                st.success(f"✅ Giữ lại {len(cols_to_keep)} cột")
                st.session_state.cleaning_log.append(f"Relevance: Giữ {len(cols_to_keep)}/{len(st.session_state.data.columns)} cột")
        
        # Step 2: Duplicates
        with st.expander("2️⃣ Remove Duplicates - Xóa dữ liệu trùng lặp", expanded=False):
            dup_count = df.duplicated().sum()
            if dup_count > 0:
                st.warning(f"⚠️ Phát hiện {dup_count} dòng trùng lặp")
                if st.button("Xóa duplicates"):
                    df = df.drop_duplicates()
                    st.success(f"✅ Đã xóa {dup_count} dòng")
                    st.session_state.cleaning_log.append(f"Duplicates: Xóa {dup_count} dòng trùng")
            else:
                st.success("✅ Không có dữ liệu trùng lặp")
        
        # Step 3: Convention
        with st.expander("3️⃣ Convention - Kiểm tra chính tả, ngữ pháp", expanded=False):
            st.info("🔍 Kiểm tra giá trị unique trong các categorical columns")
            
            for col in st.session_state.categorical_cols:
                if col in df.columns:
                    unique_vals = df[col].value_counts()
                    st.markdown(f"**{col}**: {len(unique_vals)} giá trị unique")
                    st.dataframe(unique_vals.head(10), width='stretch')
        
        # Step 4: Format
        with st.expander("4️⃣ Format - Chuyển đổi kiểu dữ liệu", expanded=False):
            
            # Convert categorical
            for col in st.session_state.categorical_cols:
                if col in df.columns:
                    df[col] = df[col].astype('category')
            
            # Convert numerical
            for col in st.session_state.numerical_cols:
                if col in df.columns:
                    try:
                        df[col] = pd.to_numeric(df[col], errors='coerce')
                    except:
                        st.warning(f"⚠️ Không thể convert {col} sang số")
            
            st.success("✅ Đã chuẩn hóa kiểu dữ liệu")
            st.dataframe(df.dtypes, width='stretch')
        
        # Step 5: Missing Values
        with st.expander("5️⃣ Missing Values - Xử lý giá trị thiếu", expanded=False):
            
            missing = df.isnull().sum()
            missing_pct = (missing.astype(float) / len(df) * 100).round(2)
            missing_df = pd.DataFrame({
                'Column': missing.index,
                'Missing': missing.values,
                'Percent': missing_pct.values
            })
            missing_df = missing_df[missing_df['Missing'] > 0]
            
            if len(missing_df) > 0:
                st.warning(f"⚠️ Phát hiện {len(missing_df)} cột có missing values")
                st.dataframe(missing_df, width='stretch')
                
                st.markdown("**Phương pháp xử lý:**")
                for col in missing_df['Column']:
                    method = st.selectbox(
                        f"Xử lý {col}:",
                        ["Giữ nguyên", "Xóa dòng", "Fill mean", "Fill median", "Fill mode", "Fill 0"],
                        key=f"missing_{col}"
                    )
                    
                    if method == "Xóa dòng":
                        df = df.dropna(subset=[col])
                    elif method == "Fill mean" and col in st.session_state.numerical_cols:
                        df[col].fillna(df[col].mean(), inplace=True)
                    elif method == "Fill median" and col in st.session_state.numerical_cols:
                        df[col].fillna(df[col].median(), inplace=True)
                    elif method == "Fill mode":
                        df[col].fillna(df[col].mode()[0], inplace=True)
                    elif method == "Fill 0":
                        df[col].fillna(0, inplace=True)
                
                if st.button("✅ Áp dụng xử lý missing values"):
                    st.success("Đã xử lý missing values")
                    st.session_state.cleaning_log.append("Missing values: Đã xử lý")
            else:
                st.success("✅ Không có missing values")
        
        # Step 6: Outliers
        with st.expander("6️⃣ Outliers - Phát hiện giá trị ngoại lai", expanded=False):
            
            if st.session_state.numerical_cols:
                for col in st.session_state.numerical_cols:
                    if col in df.columns:
                        # IQR method
                        Q1 = df[col].quantile(0.25)
                        Q3 = df[col].quantile(0.75)
                        IQR = Q3 - Q1
                        lower = Q1 - 1.5 * IQR
                        upper = Q3 + 1.5 * IQR
                        
                        outliers = df[(df[col] < lower) | (df[col] > upper)]
                        
                        if len(outliers) > 0:
                            st.warning(f"⚠️ **{col}**: {len(outliers)} outliers")
                            
                            # Boxplot
                            fig = px.box(df, y=col, title=f"Boxplot - {col}")
                            st.plotly_chart(fig, use_container_width=True)
                            
                            # Outlier treatment
                            action = st.selectbox(
                                f"Xử lý outliers {col}:",
                                ["Giữ nguyên", "Xóa outliers", "Cap outliers", "Transform log"],
                                key=f"outlier_{col}"
                            )
                            
                            if action == "Xóa outliers":
                                df = df[(df[col] >= lower) & (df[col] <= upper)]
                            elif action == "Cap outliers":
                                df[col] = df[col].clip(lower, upper)
                            elif action == "Transform log":
                                df[col] = np.log1p(df[col])
        
        # Step 7: Scaling
        with st.expander("7️⃣ Scaling - Chuẩn hóa dữ liệu số", expanded=False):
            
            if st.session_state.numerical_cols:
                st.info("📊 Scaling giúp đưa các cột số về cùng một thang đo, quan trọng cho Machine Learning")
                
                scaling_method = st.selectbox(
                    "Chọn phương pháp scaling:",
                    ["Không scaling", "StandardScaler (z-score)", "MinMaxScaler (0-1)", "RobustScaler (IQR)"]
                )
                
                if scaling_method != "Không scaling":
                    cols_to_scale = st.multiselect(
                        "Chọn columns cần scaling:",
                        st.session_state.numerical_cols,
                        default=st.session_state.numerical_cols
                    )
                    
                    if st.button("✅ Áp dụng scaling"):
                        from sklearn.preprocessing import StandardScaler, MinMaxScaler, RobustScaler
                        
                        if scaling_method == "StandardScaler (z-score)":
                            scaler = StandardScaler()
                        elif scaling_method == "MinMaxScaler (0-1)":
                            scaler = MinMaxScaler()
                        else:  # RobustScaler
                            scaler = RobustScaler()
                        
                        for col in cols_to_scale:
                            if col in df.columns:
                                df[col] = scaler.fit_transform(df[[col]])
                        
                        st.success(f"✅ Đã scaling {len(cols_to_scale)} columns với {scaling_method}")
                        st.session_state.cleaning_log.append(f"Scaling: {scaling_method} cho {len(cols_to_scale)} cột")
                        
                        # Show before/after comparison
                        st.markdown("**Trước và sau scaling:**")
                        for col in cols_to_scale[:3]:  # Show first 3
                            if col in df.columns:
                                col1, col2 = st.columns(2)
                                with col1:
                                    st.metric(f"{col} - Min", f"{st.session_state.data[col].min():.2f}")
                                    st.metric(f"{col} - Max", f"{st.session_state.data[col].max():.2f}")
                                with col2:
                                    st.metric(f"{col} - Min (scaled)", f"{df[col].min():.2f}")
                                    st.metric(f"{col} - Max (scaled)", f"{df[col].max():.2f}")
            else:
                st.info("Không có numerical columns để scaling")
        
        # Step 8: Validation
        with st.expander("8️⃣ Validation - Kiểm tra kết quả cleaning", expanded=True):
            
            col1, col2 = st.columns(2)
            
            with col1:
                st.markdown("**Before Cleaning:**")
                st.metric("Rows", len(st.session_state.data))
                st.metric("Columns", len(st.session_state.data.columns))
                st.metric("Missing", st.session_state.data.isnull().sum().sum())
            
            with col2:
                st.markdown("**After Cleaning:**")
                st.metric("Rows", len(df), delta=len(df) - len(st.session_state.data))
                st.metric("Columns", len(df.columns), delta=len(df.columns) - len(st.session_state.data.columns))
                st.metric("Missing", df.isnull().sum().sum(), delta=df.isnull().sum().sum() - st.session_state.data.isnull().sum().sum())
            
            # Cleaning log
            st.markdown("**📋 Cleaning Log:**")
            for log in st.session_state.cleaning_log:
                st.text(f"✓ {log}")
            
            if st.button("💾 Lưu cleaned data", type="primary"):
                # Convert categorical columns to object to avoid PyArrow serialization errors
                df_to_save = df.copy()
                for col in df_to_save.select_dtypes(include=['category']).columns:
                    df_to_save[col] = df_to_save[col].astype(str)
                
                st.session_state.cleaned_data = df_to_save
                st.success("✅ Đã lưu cleaned data! Chuyển sang Step 3 để phân tích EDA.")

elif workflow_step == "📊 Step 3: EDA & Insights":
    st.markdown("""
    <div class="step-card">
        <h2>📊 Bước 3: EDA & Insights</h2>
        <p>Khám phá dữ liệu và tạo insights với AI</p>
    </div>
    """, unsafe_allow_html=True)
    
    if st.session_state.cleaned_data is None:
        st.warning("⚠️ Vui lòng hoàn thành Step 2 trước!")
    else:
        df = st.session_state.cleaned_data.copy()
        
        eda_tabs = st.tabs([
            "📈 Univariate", 
            "🔗 Bivariate",
            "🎯 Multivariate",
            "🤖 AI Insights"
        ])
        
        with eda_tabs[0]:  # Univariate
            st.markdown("### 📈 Phân tích đơn biến (Univariate)")
            
            # Numerical distributions
            if st.session_state.numerical_cols:
                st.markdown("#### Phân phối Numerical Columns")
                for col in st.session_state.numerical_cols:
                    if col in df.columns:
                        fig = px.histogram(df, x=col, title=f"Distribution: {col}", 
                                         marginal="box", nbins=30)
                        st.plotly_chart(fig, use_container_width=True)
            
            # Categorical frequencies
            if st.session_state.categorical_cols:
                st.markdown("#### Tần suất Categorical Columns")
                for col in st.session_state.categorical_cols:
                    if col in df.columns:
                        value_counts = df[col].value_counts().head(10)
                        fig = px.bar(x=value_counts.index, y=value_counts.values,
                                   title=f"Frequency: {col}",
                                   labels={'x': col, 'y': 'Count'})
                        st.plotly_chart(fig, use_container_width=True)
            
            # DateTime time series
            if st.session_state.get('date_cols', []):
                st.markdown("#### Phân tích Time Series")
                for date_col in st.session_state.get('date_cols', []):
                    if date_col in df.columns and pd.api.types.is_datetime64_any_dtype(df[date_col]):
                        df_sorted = df.sort_values(date_col)
                        fig = px.line(df_sorted, x=date_col, y=df_sorted.index,
                                    title=f"Timeline: {date_col}")
                        st.plotly_chart(fig, use_container_width=True)
        
        with eda_tabs[1]:  # Bivariate
            st.markdown("### 🔗 Phân tích song biến (Bivariate)")
            
            # Correlation heatmap
            if len(st.session_state.numerical_cols) >= 2:
                numerical_data = df[st.session_state.numerical_cols].select_dtypes(include=[np.number])
                corr = numerical_data.corr()
                
                fig = px.imshow(corr, 
                              text_auto=True,
                              title="Correlation Heatmap",
                              color_continuous_scale="RdBu_r")
                st.plotly_chart(fig, use_container_width=True)
            
            # Scatter plots
            if len(st.session_state.numerical_cols) >= 2:
                st.markdown("#### Scatter Plots")
                col_x = st.selectbox("Chọn X axis:", st.session_state.numerical_cols)
                col_y = st.selectbox("Chọn Y axis:", 
                                   [c for c in st.session_state.numerical_cols if c != col_x])
                
                fig = px.scatter(df, x=col_x, y=col_y, title=f"{col_x} vs {col_y}",
                               trendline="ols")
                st.plotly_chart(fig, use_container_width=True)
            
            # Time series with numerical values
            if st.session_state.get('date_cols', []) and st.session_state.numerical_cols:
                st.markdown("#### Time Series Analysis")
                date_col = st.selectbox("Chọn Date column:", st.session_state.get('date_cols', []))
                num_col = st.selectbox("Chọn Numerical column:", st.session_state.numerical_cols)
                
                if date_col in df.columns and num_col in df.columns:
                    df_sorted = df[[date_col, num_col]].dropna().sort_values(date_col)
                    fig = px.line(df_sorted, x=date_col, y=num_col,
                                title=f"{num_col} theo thời gian")
                    st.plotly_chart(fig, use_container_width=True)
        
        with eda_tabs[2]:  # Multivariate
            st.markdown("### 🎯 Phân tích đa biến (Multivariate)")
            
            # Pairplot-style
            if len(st.session_state.numerical_cols) >= 2:
                selected_cols = st.multiselect(
                    "Chọn columns để phân tích:",
                    st.session_state.numerical_cols,
                    default=st.session_state.numerical_cols[:min(3, len(st.session_state.numerical_cols))]
                )
                
                if len(selected_cols) >= 2:
                    fig = px.scatter_matrix(df[selected_cols],
                                          title="Scatter Matrix")
                    st.plotly_chart(fig, use_container_width=True)
        
        with eda_tabs[3]:  # AI Insights
            st.markdown("### 🤖 AI-Powered Insights")
            
            client = init_gemini()
            if client:
                if st.button("🚀 Tạo AI Insights", type="primary"):
                    with st.spinner("Gemini đang phân tích dữ liệu..."):
                        # Prepare data summary
                        summary = f"""
Dataset: {st.session_state.dataset_description}
Mục tiêu: {st.session_state.analysis_goal}

Thông tin dữ liệu sau cleaning:
- Rows: {len(df)}
- Columns: {len(df.columns)}
- Date/Time: {st.session_state.date_cols}
- Categorical: {st.session_state.categorical_cols}
- Numerical: {st.session_state.numerical_cols}

Thống kê mô tả:
{df[st.session_state.numerical_cols].describe().to_string() if st.session_state.numerical_cols else 'N/A'}

Dữ liệu mẫu:
{df.head(5).to_string()}
"""
                        
                        prompt = f"""
Thực hiện **EXPLORATORY DATA ANALYSIS & DATA TRANSFORMATION** theo best practices AWS/Google/Microsoft cho dataset:

{summary}

**YÊU CẦU PHÂN TÍCH (ETL Standards - AWS/Google/Microsoft):**

## 📊 PART 1: KEY FINDINGS & PATTERNS

**1.1 Statistical Insights:**
- Phân phối numerical columns: Normal? Skewed? Có outliers?
- Top insights từ correlation matrix (nếu có)
- Xu hướng thời gian (time-series patterns nếu có date columns)

**1.2 Data Patterns:**
- Categorical distributions: Cột nào unbalanced? Có missing categories?
- Relationships: Correlations mạnh nhất? Dependencies giữa columns?
- Anomalies: Outliers, unusual patterns, data quality issues còn sót

## 🔧 PART 2: DATA TRANSFORMATION RECOMMENDATIONS (ETL Best Practices)

**2.1 Feature Engineering Suggestions:**
- **Time-based features** (nếu có date): year, month, quarter, day_of_week, is_weekend, days_since_X
- **Derived metrics**: Tính toán KPIs từ raw data (VD: total_revenue = price × quantity, profit_margin)
- **Categorical encoding**: One-hot? Label? Target encoding? (recommend cho từng column)
- **Binning**: Age groups, revenue ranges theo business logic

**2.2 Aggregation Strategies:**
- **Time-series aggregation**: Daily/Weekly/Monthly summaries
- **Categorical aggregation**: Group by product/region/segment
- **KPIs calculation**: Common metrics cho domain này (Sales: AOV, Conversion Rate; Marketing: ROAS, CTR, CPA)

**2.3 Visualization-Ready Datasets:**
- Recommend cấu trúc data cho từng chart type:
  - Line charts: [Date, Metric1, Metric2]
  - Bar charts: [Category, Value, Percentage]
  - Heatmaps: Correlation matrix pre-computed

## 💡 PART 3: BUSINESS RECOMMENDATIONS

**3.1 Actionable Insights:**
- Top 3-5 phát hiện quan trọng nhất cho business
- Cơ hội: Trends tốt để capitalize
- Risks: Vấn đề cần giải quyết ngay

**3.2 Dashboard Readiness:**
- Data đã sẵn sàng cho dashboard chưa?
- Cần transform/aggregate thêm gì?
- Recommend primary KPIs cho dashboard

## 🎯 PART 4: NEXT STEPS

**Ưu tiên cao:**
1. Feature engineering steps cụ thể
2. KPIs cần calculate
3. Data transformations cần thiết

**Ưu tiên trung bình:**
- Deeper analysis opportunities
- A/B testing ideas (nếu applicable)

---

**OUTPUT FORMAT:**
- Markdown professional, dễ đọc
- Tiếng Việt
- Có số liệu cụ thể (không chung chung)
- Actionable recommendations (KHÔNG phải "nên xem xét", phải cụ thể "Cần tính metric X = formula Y")

**DATA TRANSFORMATION PRINCIPLES:**
- Apply AWS/Google/Microsoft ETL best practices
- Prepare data for scalable dashboard
- Optimize for performance (<5MB datasets, pre-aggregated KPIs)
"""
                        success, insight = generate_ai_insight(client, prompt, temperature=0.5, max_tokens=8192)
                        if success:
                            st.session_state.eda_results['ai_insights'] = insight
                            with st.expander("📋 Full EDA & Transformation Report", expanded=True):
                                st.markdown(insight)
                        else:
                            st.error(f"❌ {insight}")
                            st.warning("💡 Vui lòng kiểm tra GEMINI_API_KEY hoặc thử lại sau.")
            else:
                st.info("💡 Cần GEMINI_API_KEY để sử dụng AI Insights")

elif workflow_step == "🎨 Step 4: Dashboard Blueprint":
    st.markdown("""
    <div class="step-card">
        <h2>🎨 Bước 4: Dashboard Blueprint</h2>
        <p>AI Expert tự động thiết kế dashboard với OQMLB framework</p>
    </div>
    """, unsafe_allow_html=True)
    
    if st.session_state.cleaned_data is None:
        st.warning("⚠️ Vui lòng hoàn thành Step 3 trước!")
    else:
        # Initialize OQMLB session state keys
        if 'oqmlb_objectives' not in st.session_state:
            st.session_state.oqmlb_objectives = None
        if 'oqmlb_questions' not in st.session_state:
            st.session_state.oqmlb_questions = None
        if 'oqmlb_metrics' not in st.session_state:
            st.session_state.oqmlb_metrics = None
        if 'oqmlb_layout' not in st.session_state:
            st.session_state.oqmlb_layout = None
        if 'oqmlb_build' not in st.session_state:
            st.session_state.oqmlb_build = None
        
        client = init_gemini()
        
        if client:
            # Main AI Generation Button
            st.markdown("### 🤖 AI Expert Dashboard Designer")
            st.info("💡 **AI Expert Mode**: Hệ thống AI sẽ đóng vai trò chuyên gia Dashboard Design để tự động tạo toàn bộ OQMLB Blueprint chuyên nghiệp. Bạn chỉ cần review và chỉnh sửa nếu muốn.")
            
            col1, col2 = st.columns([3, 1])
            with col1:
                if st.button("🚀 AI Tạo Dashboard Blueprint Tự Động (Expert Mode)", type="primary", width='stretch'):
                    try:
                        with st.spinner("🎯 AI Expert đang phân tích dữ liệu và thiết kế dashboard chuyên nghiệp..."):
                            df = st.session_state.cleaned_data
                            
                            # Comprehensive Expert Prompt
                            expert_prompt = f"""
Bạn là Senior Dashboard Design Consultant với 15+ năm kinh nghiệm trong Business Intelligence và Data Visualization. 
Nhiệm vụ của bạn là thiết kế một Dashboard Blueprint hoàn chỉnh, chuyên nghiệp theo framework OQMLB.

**THÔNG TIN DỮ LIỆU:**
- Dataset: {st.session_state.dataset_description}
- Mục tiêu phân tích: {st.session_state.analysis_goal}
- Số lượng records: {len(df)}
- Columns: {df.columns.tolist()}
- Numerical columns: {st.session_state.numerical_cols}
- Categorical columns: {st.session_state.categorical_cols}
- Date/Time columns: {st.session_state.date_cols}

**Dữ liệu mẫu:**
{df.head(10).to_string()}

**YÊU CẦU THIẾT KẾ OQMLB FRAMEWORK:**

Hãy tạo một Dashboard Blueprint hoàn chỉnh theo cấu trúc sau (bắt buộc phải có đầy đủ 5 phần):

## O - OBJECTIVES (Mục tiêu)
- Dashboard purpose: Mục đích chính của dashboard này
- Target audience: Đối tượng sử dụng (role, level)
- Key business goals: 3-5 mục tiêu kinh doanh cốt lõi
- Success criteria: Tiêu chí đánh giá thành công

## Q - QUESTIONS (Câu hỏi cần trả lời)
Chia thành 5 nhóm câu hỏi (mỗi nhóm 3-5 câu):
1. **Tổng quan (Overview)**: Câu hỏi về tình hình chung
2. **Hiệu suất (Performance)**: Đo lường KPIs và metrics
3. **Xu hướng (Trends)**: Patterns theo thời gian
4. **So sánh (Comparison)**: So sánh giữa các segments
5. **Hành động (Action)**: Insights để ra quyết định

## M - METRICS (Chỉ số đo lường)

**⭐ PRIMARY KPIs** (3-5 chỉ số quan trọng nhất - theo thứ tự ưu tiên):
Với mỗi Primary KPI, cung cấp:
- **KPI Name**: Tên metric
- **Priority**: [P1/P2/P3] - P1 là quan trọng nhất (Hero KPI), P2-P3 supporting
- **Formula**: Công thức tính (nếu cần)
- **Target/Threshold**: Mục tiêu hoặc ngưỡng cảnh báo cụ thể
- **Business Question Answered**: Link đến câu hỏi nào trong phần Q
- **Action Recommendation**: 
  * Nếu metric > threshold → Nên làm gì
  * Nếu metric < threshold → Nên làm gì
  * Xu hướng tốt/xấu là gì

**Secondary Metrics** (metrics hỗ trợ):
- Liệt kê với công thức

**Calculated Metrics** (cần tính toán):
- Công thức chi tiết để implement

## L - LAYOUT (Bố cục Dashboard)

**🎨 DESIGN PRINCIPLES (BẮT BUỘC):**
- **F-Pattern Layout**: Hero KPI (P1) phải ở top-left, người dùng đọc từ trái→phải, trên→dưới
- **Visual Hierarchy**: Metrics quan trọng hơn → kích thước lớn hơn, vị trí nổi bật hơn
- **Clarity First**: Limit 6-8 visualizations per page, tránh overcrowd
- **Whitespace**: Khoảng trắng giữa components để dễ đọc
- **Color Consistency**: Tối đa 3-4 màu chính, sử dụng nhất quán

**📱 PAGE STRUCTURE** (thiết kế 3-4 pages/tabs):

**Page [X]: [Tên trang]**
- **Purpose**: Mục đích và target users
- **Hero Section** (Top-Left, most prominent):
  * Primary KPI (P1): [Tên KPI] - Large metric card với trend indicator
- **Supporting KPIs** (Top-Right, secondary prominence):
  * 2-3 P2/P3 KPIs - Medium size cards
- **Visualizations** (Main content area, limit 4-6 charts):
  * Chart 1: [Type] - [Data] - [Insight] - [Position: left/right/center]
  * Chart 2: [Type] - [Data] - [Insight] - [Position: ...]
  * (Max 6 charts - prioritize quality over quantity)
- **Contextual Filters** (Sidebar or top):
  * Filter 1: [Category/Time range] - Default value
  * Filter 2: [Segment/Region] - Options
- **AI Insights Panel** (Right sidebar or bottom):
  * Key findings từ data
  * Action recommendations based on current metrics
- **Layout Grid**: 
  * Row 1: Hero KPI (col-span-2) + Supporting KPIs (col-span-1 each)
  * Row 2-3: Main visualizations (balanced grid)
  * Row 4: Details/drill-down section

## B - BUILD (Hướng dẫn triển khai)
- **Data preparation**: Xử lý data cần thiết
- **Technology stack**: Tools/platforms đề xuất (Streamlit, PowerBI, Tableau, etc.)
- **Implementation priority**: Phases triển khai (Phase 1, 2, 3)
- **Best practices**: Guidelines cho developers
- **Performance considerations**: Optimize queries, caching
- **Testing checklist**: Các điểm cần test

**YÊU CẦU OUTPUT (BẮT BUỘC):**
- Trả lời bằng tiếng Việt
- Professional, chi tiết, actionable
- Format Markdown đẹp với headers, bullet points, tables nếu cần
- Đưa ra recommendations cụ thể dựa trên data thực tế
- Apply best practices trong Dashboard Design (Gestalt principles, color theory, data-ink ratio)

**✨ 5 TIÊU CHÍ CHẤT LƯỢNG DASHBOARD (phải đáp ứng):**
1. ❶ **Informative**: Mỗi chart/KPI phải trả lời rõ ràng business question nào
2. ❷ **Clarity**: Dễ hiểu trong 5-10 giây, chart types phù hợp, hierarchy rõ ràng
3. ❸ **Design**: F-pattern layout, whitespace hợp lý, color consistency
4. ❹ **Interactivity**: Filters có context, drill-down paths rõ ràng
5. ❺ **Actionable Insights**: Mỗi metric phải có action recommendations - "Nếu X thì làm Y"

**FORMAT TEMPLATE BẮT BUỘC (KHÔNG ĐƯỢC THAY ĐỔI HEADERS):**
Bạn PHẢI sử dụng chính xác các headers sau (bao gồm cả dấu ##):

## O - OBJECTIVES
[Nội dung về objectives...]

## Q - QUESTIONS
[Nội dung về questions...]

## M - METRICS
[Nội dung về metrics...]

## L - LAYOUT
[Nội dung về layout...]

## B - BUILD
[Nội dung về build...]

QUAN TRỌNG: Phải có đầy đủ 5 sections với headers chính xác như trên.
                            """
                            
                            # Generate complete OQMLB (needs more tokens for comprehensive blueprint)
                            success, full_blueprint = generate_ai_insight(client, expert_prompt, temperature=0.5, max_tokens=8192)
                            
                            if not success:
                                st.error(f"❌ {full_blueprint}")
                                st.warning("💡 Vui lòng kiểm tra GEMINI_API_KEY hoặc thử lại sau.")
                                st.stop()
                            
                            # Parse sections with improved logic
                            sections = full_blueprint.split('## ')
                            
                            # Reset sections first
                            st.session_state.oqmlb_objectives = None
                            st.session_state.oqmlb_questions = None
                            st.session_state.oqmlb_metrics = None
                            st.session_state.oqmlb_layout = None
                            st.session_state.oqmlb_build = None
                            
                            # Debug: Show what sections were found
                            section_headers = [s.split('\n')[0][:50] for s in sections if s.strip()]
                            
                            for section in sections:
                                section_text = section.strip()
                                if section_text:
                                    section_upper = section_text.upper()
                                    
                                    # More flexible parsing - check if starts with O/Q/M/L/B
                                    if section_upper.startswith('O -') or section_upper.startswith('OBJECTIVES'):
                                        st.session_state.oqmlb_objectives = '## ' + section_text
                                    elif section_upper.startswith('Q -') or section_upper.startswith('QUESTIONS'):
                                        st.session_state.oqmlb_questions = '## ' + section_text
                                    elif section_upper.startswith('M -') or section_upper.startswith('METRICS'):
                                        st.session_state.oqmlb_metrics = '## ' + section_text
                                    elif section_upper.startswith('L -') or section_upper.startswith('LAYOUT'):
                                        st.session_state.oqmlb_layout = '## ' + section_text
                                    elif section_upper.startswith('B -') or section_upper.startswith('BUILD'):
                                        st.session_state.oqmlb_build = '## ' + section_text
                            
                            # Store complete blueprint
                            st.session_state.dashboard_blueprint = full_blueprint
                            
                            # Validate all sections are present
                            missing_sections = []
                            if not st.session_state.oqmlb_objectives:
                                missing_sections.append("Objectives")
                            if not st.session_state.oqmlb_questions:
                                missing_sections.append("Questions")
                            if not st.session_state.oqmlb_metrics:
                                missing_sections.append("Metrics")
                            if not st.session_state.oqmlb_layout:
                                missing_sections.append("Layout")
                            if not st.session_state.oqmlb_build:
                                missing_sections.append("Build")
                            
                            if missing_sections:
                                st.warning(f"⚠️ Một số sections chưa được parse đúng: {', '.join(missing_sections)}.")
                                st.info(f"📋 Headers tìm thấy: {', '.join(section_headers[:10])}")
                                st.info("💡 Xem toàn bộ blueprint ở tab 'Complete Blueprint' bên dưới.")
                            else:
                                st.success("✅ AI Expert đã tạo xong Dashboard Blueprint chuyên nghiệp!")
                            
                            # Parse structured data for Step 5 dashboard generation
                            try:
                                structured_data = parse_blueprint_structured_data(
                                    st.session_state.oqmlb_metrics or "",
                                    st.session_state.oqmlb_layout or ""
                                )
                                st.session_state.blueprint_structured_data = structured_data
                                
                                # Debug info
                                if structured_data.get('hero_kpi'):
                                    st.info(f"🎯 Hero KPI detected: {structured_data['hero_kpi'].get('name', 'N/A')}")
                            except Exception as e:
                                st.warning(f"⚠️ Không thể parse structured data: {str(e)}")
                                st.session_state.blueprint_structured_data = None
                            
                            st.rerun()
                    
                    except Exception as e:
                        st.error(f"❌ Lỗi khi tạo blueprint: {str(e)}")
                        st.error("Vui lòng thử lại hoặc kiểm tra kết nối internet và GEMINI_API_KEY.")
            
            with col2:
                if st.button("🔄 Reset", width='stretch'):
                    st.session_state.oqmlb_objectives = None
                    st.session_state.oqmlb_questions = None
                    st.session_state.oqmlb_metrics = None
                    st.session_state.oqmlb_layout = None
                    st.session_state.oqmlb_build = None
                    st.session_state.dashboard_blueprint = None
                    st.rerun()
            
            st.markdown("---")
            
            # Display OQMLB Framework in tabs
            if st.session_state.oqmlb_objectives:
                st.markdown("### 📋 OQMLB Framework (AI Generated)")
                
                framework_tabs = st.tabs([
                    "🎯 Objectives",
                    "❓ Questions", 
                    "📊 Metrics",
                    "📐 Layout",
                    "🔨 Build",
                    "📥 Complete Blueprint"
                ])
                
                with framework_tabs[0]:  # Objectives
                    st.markdown(st.session_state.oqmlb_objectives)
                    
                    with st.expander("✏️ Chỉnh sửa Objectives"):
                        edited_obj = st.text_area(
                            "Edit nội dung:",
                            value=st.session_state.oqmlb_objectives,
                            height=300,
                            key="edit_objectives"
                        )
                        if st.button("💾 Lưu thay đổi", key="save_obj"):
                            st.session_state.oqmlb_objectives = edited_obj
                            st.success("Đã lưu!")
                
                with framework_tabs[1]:  # Questions
                    st.markdown(st.session_state.oqmlb_questions if st.session_state.oqmlb_questions else "_Chưa có dữ liệu_")
                    
                    with st.expander("✏️ Chỉnh sửa Questions"):
                        edited_q = st.text_area(
                            "Edit nội dung:",
                            value=st.session_state.oqmlb_questions or "",
                            height=300,
                            key="edit_questions"
                        )
                        if st.button("💾 Lưu thay đổi", key="save_q"):
                            st.session_state.oqmlb_questions = edited_q
                            st.success("Đã lưu!")
                
                with framework_tabs[2]:  # Metrics
                    st.markdown(st.session_state.oqmlb_metrics if st.session_state.oqmlb_metrics else "_Chưa có dữ liệu_")
                    
                    with st.expander("✏️ Chỉnh sửa Metrics"):
                        edited_m = st.text_area(
                            "Edit nội dung:",
                            value=st.session_state.oqmlb_metrics or "",
                            height=300,
                            key="edit_metrics"
                        )
                        if st.button("💾 Lưu thay đổi", key="save_m"):
                            st.session_state.oqmlb_metrics = edited_m
                            st.success("Đã lưu!")
                
                with framework_tabs[3]:  # Layout
                    st.markdown(st.session_state.oqmlb_layout if st.session_state.oqmlb_layout else "_Chưa có dữ liệu_")
                    
                    with st.expander("✏️ Chỉnh sửa Layout"):
                        edited_l = st.text_area(
                            "Edit nội dung:",
                            value=st.session_state.oqmlb_layout or "",
                            height=300,
                            key="edit_layout"
                        )
                        if st.button("💾 Lưu thay đổi", key="save_l"):
                            st.session_state.oqmlb_layout = edited_l
                            st.success("Đã lưu!")
                
                with framework_tabs[4]:  # Build
                    st.markdown(st.session_state.oqmlb_build if st.session_state.oqmlb_build else "_Chưa có dữ liệu_")
                    
                    with st.expander("✏️ Chỉnh sửa Build"):
                        edited_b = st.text_area(
                            "Edit nội dung:",
                            value=st.session_state.oqmlb_build or "",
                            height=300,
                            key="edit_build"
                        )
                        if st.button("💾 Lưu thay đổi", key="save_b"):
                            st.session_state.oqmlb_build = edited_b
                            st.success("Đã lưu!")
                
                with framework_tabs[5]:  # Complete Blueprint
                    st.markdown("### 📋 Dashboard Blueprint Hoàn Chỉnh")
                    
                    if st.session_state.dashboard_blueprint:
                        st.markdown(st.session_state.dashboard_blueprint)
                        
                        # Download options
                        st.markdown("---")
                        st.markdown("### 📥 Tải xuống Blueprint")
                        
                        col1, col2 = st.columns(2)
                        with col1:
                            st.download_button(
                                label="📥 Download Blueprint (.md)",
                                data=st.session_state.dashboard_blueprint,
                                file_name=f"dashboard_blueprint_{datetime.now().strftime('%Y%m%d_%H%M')}.md",
                                mime="text/markdown",
                                width='stretch'
                            )
                        
                        with col2:
                            # Also offer TXT format
                            st.download_button(
                                label="📄 Download Blueprint (.txt)",
                                data=st.session_state.dashboard_blueprint,
                                file_name=f"dashboard_blueprint_{datetime.now().strftime('%Y%m%d_%H%M')}.txt",
                                mime="text/plain",
                                width='stretch'
                            )
            else:
                st.info("👆 Click nút **'🚀 AI Tạo Dashboard Blueprint Tự Động'** để AI Expert tạo blueprint chuyên nghiệp cho bạn!")
        
        else:
            st.error("❌ Cần GEMINI_API_KEY để sử dụng AI Expert Dashboard Designer")
            st.info("💡 Vui lòng cấu hình GEMINI_API_KEY trong Secrets để kích hoạt tính năng này.")

elif workflow_step == "🚀 Step 5: Build Dashboard":
    st.markdown("""
    <div class="step-card">
        <h2>🚀 Bước 5: Auto-Generated Interactive Dashboard</h2>
        <p>✨ Dashboard được tạo tự động từ AI Blueprint của bạn</p>
    </div>
    """, unsafe_allow_html=True)
    
    if st.session_state.cleaned_data is None:
        st.warning("⚠️ Vui lòng hoàn thành Step 2 (Data Cleaning) trước!")
    else:
        df = st.session_state.cleaned_data
        
        # Success message
        st.success("✅ Dashboard đã được tạo tự động từ OQMLB Blueprint!")
        
        # Dashboard tabs based on common structure
        dashboard_tabs = st.tabs([
            "📊 Overview Dashboard",
            "📈 Detailed Analysis", 
            "🎯 KPIs & Metrics",
            "📥 Export & Code"
        ])
        
        with dashboard_tabs[0]:  # Overview Dashboard
            st.markdown("### 📊 Executive Dashboard")
            st.markdown("*Tự động tạo theo OQMLB Blueprint + 5 Tiêu Chí Chuyên Nghiệp*")
            
            # AI Dashboard Insights (Bricks.ai-inspired feature)
            st.markdown("---")
            client = init_gemini()
            if client:
                # Handle regeneration or first-time generation
                should_generate = False
                is_regenerating = False
                
                if st.session_state.regenerate_insights_flag:
                    st.session_state.regenerate_insights_flag = False
                    should_generate = True
                    is_regenerating = True
                
                # Show UI based on current state
                if not st.session_state.dashboard_insights and not should_generate:
                    # No insights yet - show generate button
                    if st.button("🤖 Generate AI Dashboard Insights", type="primary", use_container_width=True, help="AI phân tích toàn bộ dashboard và tạo insights với specific numbers"):
                        should_generate = True
                    
                    st.info("💡 Click **'Generate AI Dashboard Insights'** để AI phân tích toàn bộ dashboard với specific numbers (giống Bricks.ai)")
                
                elif st.session_state.dashboard_insights and not should_generate:
                    # Insights available - show with regenerate option
                    col_success, col_regen = st.columns([3, 2])
                    with col_success:
                        st.success("✅ AI Dashboard Insights đã sẵn sàng!")
                    with col_regen:
                        if st.button("🔄 Regenerate Insights", use_container_width=True):
                            st.session_state.regenerate_insights_flag = True
                            st.rerun()
                    
                    with st.expander("💡 AI Dashboard Insights (Bricks.ai Style)", expanded=True):
                        st.markdown(st.session_state.dashboard_insights)
                
                # Execute generation if triggered
                if should_generate:
                    spinner_msg = "🔄 AI đang regenerate insights..." if is_regenerating else "🤖 AI đang phân tích dashboard data..."
                    with st.spinner(spinner_msg):
                        success, insights = generate_dashboard_insights(
                            client,
                            df,
                            st.session_state.numerical_cols or [],
                            st.session_state.categorical_cols or []
                        )
                        if success:
                            st.session_state.dashboard_insights = insights
                            st.rerun()
                        else:
                            st.error(f"❌ Lỗi: {insights}")
                            if is_regenerating:
                                st.warning("💡 Regeneration failed, keeping old insights")
                            st.stop()
                
                st.markdown("---")
                
                # Per-Chart Insights (Bricks.ai per-visual insights)
                st.markdown("#### 💡 Per-Chart Insights (Bricks.ai Style)")
                st.caption("*Insights cho TỪNG chart - giống Bricks.ai professional reports*")
                
                # Handle regeneration for per-chart insights
                should_generate_charts = False
                is_regenerating_charts = False
                
                if st.session_state.regenerate_chart_insights_flag:
                    st.session_state.regenerate_chart_insights_flag = False
                    should_generate_charts = True
                    is_regenerating_charts = True
                
                # Show UI based on current state
                if not st.session_state.per_chart_insights and not should_generate_charts:
                    col_gen_chart = st.columns(1)[0]
                    if st.button("✨ Generate Per-Chart Insights", type="secondary", use_container_width=True, help="AI tạo insights riêng cho TỪNG chart (2-3 bullets per visual)"):
                        should_generate_charts = True
                    
                    st.info("✨ Click để tạo insights box bên cạnh mỗi chart (Bricks.ai format)")
                
                elif st.session_state.per_chart_insights and not should_generate_charts:
                    col_chart_success, col_chart_regen = st.columns([3, 2])
                    with col_chart_success:
                        st.success(f"✅ Per-Chart Insights generated cho {len(st.session_state.per_chart_insights)} charts!")
                    with col_chart_regen:
                        if st.button("🔄 Regenerate Per-Chart", use_container_width=True):
                            st.session_state.regenerate_chart_insights_flag = True
                            st.rerun()
                
                # Execute generation if triggered
                if should_generate_charts:
                    spinner_msg = "🔄 AI đang regenerate per-chart insights..." if is_regenerating_charts else "✨ AI đang tạo insights cho từng chart..."
                    with st.spinner(spinner_msg):
                        success, insights_dict = generate_per_chart_insights(
                            client,
                            df,
                            st.session_state.numerical_cols or [],
                            st.session_state.categorical_cols or []
                        )
                        if success:
                            st.session_state.per_chart_insights = insights_dict
                            st.rerun()
                        else:
                            st.error("❌ Lỗi khi generate per-chart insights")
                            if is_regenerating_charts:
                                st.warning("💡 Regeneration failed, keeping old insights")
                            st.stop()
            
            st.markdown("---")
            
            # 🎯 PHASE 1: INTERACTIVE FILTER PANEL (Better than Bricks.ai click-to-filter!)
            st.markdown("### 🔍 Interactive Data Filters")
            st.caption("*Filter toàn bộ dashboard - tất cả charts sẽ update real-time*")
            
            # Create filter UI
            filter_col1, filter_col2, filter_col3 = st.columns([3, 3, 1])
            
            with filter_col1:
                # Categorical filters
                if st.session_state.categorical_cols:
                    selected_cat_col = st.selectbox(
                        "📊 Filter by Category",
                        options=["(No filter)"] + st.session_state.categorical_cols,
                        key="filter_cat_col"
                    )
                    
                    if selected_cat_col != "(No filter)":
                        unique_vals = df[selected_cat_col].unique().tolist()
                        selected_vals = st.multiselect(
                            f"Select {selected_cat_col} values:",
                            options=unique_vals,
                            default=st.session_state.active_filters.get(selected_cat_col, unique_vals),
                            key=f"filter_{selected_cat_col}"
                        )
                        st.session_state.active_filters[selected_cat_col] = selected_vals
                    else:
                        # Clear categorical filters
                        for col in st.session_state.categorical_cols:
                            if col in st.session_state.active_filters:
                                del st.session_state.active_filters[col]
            
            with filter_col2:
                # Numerical range filters
                if st.session_state.numerical_cols:
                    selected_num_col = st.selectbox(
                        "📈 Filter by Range",
                        options=["(No filter)"] + st.session_state.numerical_cols,
                        key="filter_num_col"
                    )
                    
                    if selected_num_col != "(No filter)":
                        col_min = float(df[selected_num_col].min())
                        col_max = float(df[selected_num_col].max())
                        default_range = st.session_state.active_filters.get(f"{selected_num_col}_range", (col_min, col_max))
                        
                        selected_range = st.slider(
                            f"{selected_num_col} range:",
                            min_value=col_min,
                            max_value=col_max,
                            value=default_range,
                            key=f"range_{selected_num_col}"
                        )
                        st.session_state.active_filters[f"{selected_num_col}_range"] = selected_range
                        st.session_state.active_filters[f"{selected_num_col}_range_col"] = selected_num_col
                    else:
                        # Clear numerical filters
                        for key in list(st.session_state.active_filters.keys()):
                            if key.endswith('_range') or key.endswith('_range_col'):
                                del st.session_state.active_filters[key]
            
            with filter_col3:
                st.write("")  # Spacing
                st.write("")  # Spacing
                if st.button("🔄 Clear All", help="Reset tất cả filters"):
                    st.session_state.active_filters = {}
                    st.rerun()
            
            # Apply filters to dataframe
            filtered_df = df.copy()
            active_filter_count = 0
            
            # Apply categorical filters
            for col in st.session_state.categorical_cols:
                if col in st.session_state.active_filters and st.session_state.active_filters[col]:
                    filtered_df = filtered_df[filtered_df[col].isin(st.session_state.active_filters[col])]
                    active_filter_count += 1
            
            # Apply numerical range filters
            for key, value in list(st.session_state.active_filters.items()):
                if key.endswith('_range_col'):
                    range_col = value
                    range_key = key.replace('_range_col', '_range')
                    if range_key in st.session_state.active_filters:
                        min_val, max_val = st.session_state.active_filters[range_key]
                        filtered_df = filtered_df[
                            (filtered_df[range_col] >= min_val) & (filtered_df[range_col] <= max_val)
                        ]
                        active_filter_count += 1
            
            # Show filter status
            if active_filter_count > 0:
                st.success(f"✅ {active_filter_count} filter(s) active | Showing {len(filtered_df):,} / {len(df):,} records ({len(filtered_df)/len(df)*100:.1f}%)")
            else:
                st.info("ℹ️ No filters applied - showing all data")
            
            # Use filtered dataframe for all visualizations below
            df = filtered_df
            
            st.markdown("---")
            
            # ❶ INFORMATIVE: Get structured data from blueprint
            blueprint_data = st.session_state.blueprint_structured_data
            
            # F-PATTERN LAYOUT: Hero KPI (top-left) + Supporting KPIs (top-right)
            if st.session_state.numerical_cols:
                st.markdown("#### 🎯 Key Performance Indicators")
                
                # Check if we have parsed blueprint KPIs
                if blueprint_data and blueprint_data.get('primary_kpis'):
                    primary_kpis = blueprint_data['primary_kpis'][:5]  # Max 5 KPIs
                    hero_kpi = blueprint_data.get('hero_kpi')  # P1 priority
                    
                    # Hero KPI Section (top-left, prominent)
                    if hero_kpi and hero_kpi.get('name'):
                        st.markdown(f"##### 🏆 {hero_kpi['name']} (Primary)")
                        
                        # Find matching column for hero KPI (with safety checks)
                        hero_col = None
                        if st.session_state.numerical_cols:
                            for col in st.session_state.numerical_cols:
                                if hero_kpi['name'].lower() in col.lower() or col.lower() in hero_kpi['name'].lower():
                                    hero_col = col
                                    break
                            
                            if not hero_col:
                                hero_col = st.session_state.numerical_cols[0]
                        
                        if hero_col:
                            hero_val = df[hero_col].sum()
                            hero_avg = df[hero_col].mean()
                            
                            # Large hero metric card
                            col_hero, col_insight = st.columns([2, 3])
                            with col_hero:
                                st.metric(
                                    label=f"{hero_col.replace('_', ' ').title()}",
                                    value=f"{hero_val:,.0f}" if hero_val > 1000 else f"{hero_val:.2f}",
                                    delta=f"Avg: {hero_avg:.1f}",
                                    help=f"⭐ Primary KPI - Priority P1"
                                )
                        
                            # ❺ ACTIONABLE INSIGHTS: Show recommendations
                            with col_insight:
                                action_rec = hero_kpi.get('action_recommendations') or ""
                                threshold = hero_kpi.get('threshold') or ""
                                
                                if action_rec:
                                    st.info(f"💡 **Action Recommendation:**\n\n{str(action_rec)[:200]}...")
                                elif threshold:
                                    st.info(f"🎯 **Target:** {str(threshold)}")
                            
                            st.markdown("---")
                    
                    # Supporting KPIs (P2, P3)
                    supporting_kpis = [kpi for kpi in primary_kpis if kpi.get('priority') in ['P2', 'P3']][:4]
                    if supporting_kpis:
                        st.markdown("##### 📊 Supporting Metrics")
                        kpi_cols = st.columns(len(supporting_kpis))
                        
                        for i, kpi in enumerate(supporting_kpis):
                            with kpi_cols[i]:
                                # Match KPI to column (with safety)
                                matched_col = None
                                if st.session_state.numerical_cols:
                                    for col in st.session_state.numerical_cols:
                                        if kpi['name'].lower() in col.lower() or col.lower() in kpi['name'].lower():
                                            matched_col = col
                                            break
                                    
                                    if not matched_col and i < len(st.session_state.numerical_cols):
                                        matched_col = st.session_state.numerical_cols[i]
                                
                                if matched_col and matched_col in df.columns:
                                    val = df[matched_col].mean()
                                    std = df[matched_col].std()
                                    
                                    # ❸ DESIGN: Color coding based on priority
                                    priority_emoji = "🥈" if kpi.get('priority') == 'P2' else "🥉"
                                    
                                    st.metric(
                                        label=f"{priority_emoji} {matched_col.replace('_', ' ').title()}",
                                        value=f"{val:,.0f}" if val > 1000 else f"{val:.2f}",
                                        delta=f"±{std:.1f}",
                                        help=f"Priority: {kpi.get('priority', 'P2')}\nTarget: {kpi.get('threshold', 'N/A')}"
                                    )
                else:
                    # Fallback: Use numerical columns directly
                    num_kpis = min(4, len(st.session_state.numerical_cols))
                    kpi_cols = st.columns(num_kpis)
                    
                    for i, col_name in enumerate(st.session_state.numerical_cols[:num_kpis]):
                        with kpi_cols[i]:
                            col_mean = df[col_name].mean()
                            col_std = df[col_name].std()
                            col_min = df[col_name].min()
                            col_max = df[col_name].max()
                            
                            st.metric(
                                label=col_name.replace('_', ' ').title(),
                                value=f"{col_mean:,.0f}" if col_mean > 1000 else f"{col_mean:.2f}",
                                delta=f"±{col_std:.1f}",
                                help=f"Min: {col_min:,.0f} | Max: {col_max:,.0f}"
                            )
            
            st.markdown("---")
            
            # ❺ ACTIONABLE INSIGHTS PANEL
            if blueprint_data and blueprint_data.get('primary_kpis'):
                with st.expander("💡 AI Insights & Recommendations", expanded=False):
                    st.markdown("### 🎯 Key Findings & Next Actions")
                    
                    for kpi in blueprint_data['primary_kpis'][:3]:
                        if kpi.get('action_recommendations'):
                            st.markdown(f"**{kpi['name']}:**")
                            st.write(kpi['action_recommendations'])
                            st.markdown("---")
            
            st.markdown("---")
            
            # ❷ CLARITY: Main Visualizations (Chart Selection Rules)
            st.markdown("### 📈 Key Visualizations")
            st.caption("*Charts tự động chọn theo best practices: Line for trends, Bar for comparisons*")
            
            if st.session_state.numerical_cols and len(st.session_state.numerical_cols) >= 2:
                # Chart 1: Trend Line Chart (Bricks.ai layout: Chart | Insights)
                st.markdown("#### 📈 Trend Analysis")
                st.caption("❓ **Answers:** How is performance changing over time?")
                
                chart_col, insights_col = st.columns([2.5, 1.5])
                
                with chart_col:
                    # LINE CHART - For showing trends over time (Edward Tufte guideline)
                    fig = px.line(
                        df.reset_index(), 
                        y=st.session_state.numerical_cols[0],
                        title=f"{st.session_state.numerical_cols[0].replace('_', ' ').title()} Trend",
                        labels={'index': 'Record #', st.session_state.numerical_cols[0]: st.session_state.numerical_cols[0].replace('_', ' ').title()}
                    )
                    # ❸ DESIGN: Semantic colors (Blue = neutral/informational)
                    fig.update_traces(line_color='#0066CC', line_width=3)
                    fig.update_layout(
                        title_font_size=18,
                        title_font_family="Arial",
                        showlegend=False,
                        margin=dict(t=50, b=50, l=50, r=20)  # Whitespace
                    )
                    st.plotly_chart(fig, use_container_width=True)
                
                with insights_col:
                    # Display insights if available
                    if st.session_state.per_chart_insights.get('trend_chart'):
                        st.markdown("""
                        <div style="background-color: #f0f8ff !important; border-left: 4px solid #0066CC !important; padding: 15px !important; border-radius: 5px !important; margin-top: 20px !important;">
                            <p style="margin: 0 !important; font-size: 14px !important; font-weight: 600 !important; color: #0066CC !important; margin-bottom: 10px !important;">💡 Key Insights</p>
                        """, unsafe_allow_html=True)
                        st.markdown(st.session_state.per_chart_insights['trend_chart'])
                        st.markdown("</div>", unsafe_allow_html=True)
                    else:
                        st.info("💡 Generate Per-Chart Insights để xem insights riêng cho chart này")
                
                st.markdown("---")
                
                # Chart 2: Distribution Histogram (Bricks.ai layout: Chart | Insights)
                st.markdown("#### 📊 Distribution Analysis")
                st.caption("❓ **Answers:** What is the value distribution?")
                
                chart_col2, insights_col2 = st.columns([2.5, 1.5])
                
                with chart_col2:
                    # HISTOGRAM - For showing distribution patterns
                    fig = px.histogram(
                        df, 
                        x=st.session_state.numerical_cols[min(1, len(st.session_state.numerical_cols)-1)],
                        title=f"{st.session_state.numerical_cols[min(1, len(st.session_state.numerical_cols)-1)].replace('_', ' ').title()} Distribution",
                        labels={st.session_state.numerical_cols[min(1, len(st.session_state.numerical_cols)-1)]: 'Value'},
                        nbins=20
                    )
                    fig.update_traces(marker_color='#667eea', marker_line_color='white', marker_line_width=1)
                    fig.update_layout(
                        title_font_size=18,
                        showlegend=False,
                        margin=dict(t=50, b=50, l=50, r=20)
                    )
                    st.plotly_chart(fig, use_container_width=True)
                
                with insights_col2:
                    # Display insights if available
                    if st.session_state.per_chart_insights.get('distribution_chart'):
                        st.markdown("""
                        <div style="background-color: #f0f8ff !important; border-left: 4px solid #667eea !important; padding: 15px !important; border-radius: 5px !important; margin-top: 20px !important;">
                            <p style="margin: 0 !important; font-size: 14px !important; font-weight: 600 !important; color: #667eea !important; margin-bottom: 10px !important;">💡 Key Insights</p>
                        """, unsafe_allow_html=True)
                        st.markdown(st.session_state.per_chart_insights['distribution_chart'])
                        st.markdown("</div>", unsafe_allow_html=True)
                    else:
                        st.info("💡 Generate Per-Chart Insights để xem insights riêng cho chart này")
            
            # ❹ INTERACTIVITY: Categorical breakdown with filters
            if st.session_state.categorical_cols and st.session_state.numerical_cols:
                st.markdown("---")
                st.markdown("### 🎨 Category Breakdown & Comparison")
                st.caption("❓ **Answers:** Which categories perform best/worst?")
                
                # Interactive filters
                col_filter1, col_filter2 = st.columns(2)
                with col_filter1:
                    cat_col = st.selectbox(
                        "📊 Select Category:",
                        st.session_state.categorical_cols,
                        key="overview_cat",
                        help="Choose dimension to analyze"
                    )
                with col_filter2:
                    num_col = st.selectbox(
                        "📈 Select Metric:",
                        st.session_state.numerical_cols,
                        key="overview_num",
                        help="Choose metric to measure"
                    )
                
                # Chart 3: Category Bar Chart (Bricks.ai layout: Chart | Insights)
                chart_col3, insights_col3 = st.columns([2.5, 1.5])
                
                with chart_col3:
                    # BAR CHART - For categorical comparisons (Top 10 only)
                    grouped = df.groupby(cat_col)[num_col].sum().sort_values(ascending=False).head(10)
                    
                    # ❸ DESIGN: Gradient color (best = green, worst = orange)
                    colors = ['#28A745' if i == 0 else '#0066CC' if i < 3 else '#6C757D' 
                             for i in range(len(grouped))]
                    
                    fig = px.bar(
                        x=grouped.index,
                        y=grouped.values,
                        title=f"Top 10: {num_col.replace('_', ' ').title()} by {cat_col.replace('_', ' ').title()}",
                        labels={'x': cat_col.replace('_', ' ').title(), 'y': num_col.replace('_', ' ').title()}
                    )
                    fig.update_traces(marker_color=colors, marker_line_color='white', marker_line_width=1.5)
                    fig.update_layout(
                        title_font_size=18,
                        showlegend=False,
                        margin=dict(t=50, b=80, l=50, r=20),
                        xaxis_tickangle=-45  # Readability
                    )
                    
                    # ❺ ACTIONABLE: Add annotation for top performer
                    if len(grouped) > 0:
                        max_val = grouped.max()
                        max_cat = grouped.idxmax()
                        fig.add_annotation(
                            x=max_cat,
                            y=max_val,
                            text=f"🏆 Best: {max_cat}",
                            showarrow=True,
                            arrowhead=2,
                            arrowcolor="#28A745",
                            font=dict(size=12, color="#28A745")
                        )
                    
                    st.plotly_chart(fig, use_container_width=True)
                    
                    # Show summary stats
                    st.caption(f"💡 **Insight:** Top performer '{grouped.idxmax()}' accounts for {(grouped.max() / grouped.sum() * 100):.1f}% of total {num_col}")
                
                with insights_col3:
                    # Display insights if available
                    if st.session_state.per_chart_insights.get('category_bar_chart'):
                        st.markdown("""
                        <div style="background-color: #f0f8ff !important; border-left: 4px solid #28A745 !important; padding: 15px !important; border-radius: 5px !important; margin-top: 20px !important;">
                            <p style="margin: 0 !important; font-size: 14px !important; font-weight: 600 !important; color: #28A745 !important; margin-bottom: 10px !important;">💡 Key Insights</p>
                        """, unsafe_allow_html=True)
                        st.markdown(st.session_state.per_chart_insights['category_bar_chart'])
                        st.markdown("</div>", unsafe_allow_html=True)
                    else:
                        st.info("💡 Generate Per-Chart Insights để xem insights riêng cho chart này")
        
        with dashboard_tabs[1]:  # Detailed Analysis
            st.markdown("### 📈 Detailed Analysis Dashboard")
            
            # Interactive filters
            st.markdown("#### 🔍 Filters")
            filter_cols = st.columns(3)
            
            active_filters = {}
            
            # Category filters
            if st.session_state.categorical_cols:
                with filter_cols[0]:
                    selected_cat = st.selectbox(
                        "Filter by Category:",
                        ['All'] + st.session_state.categorical_cols,
                        key="detail_cat_filter"
                    )
                    
                    if selected_cat != 'All':
                        cat_values = ['All'] + df[selected_cat].unique().tolist()
                        selected_val = st.selectbox(
                            f"Select {selected_cat}:",
                            cat_values,
                            key="detail_cat_val"
                        )
                        if selected_val != 'All':
                            active_filters[selected_cat] = selected_val
            
            # Apply filters
            filtered_df = df.copy()
            for col, val in active_filters.items():
                filtered_df = filtered_df[filtered_df[col] == val]
            
            st.info(f"📊 Showing {len(filtered_df)} of {len(df)} records")
            
            # Correlation heatmap (Bricks.ai layout: Chart | Insights)
            if st.session_state.numerical_cols and len(st.session_state.numerical_cols) >= 2:
                st.markdown("#### 🔥 Correlation Heatmap")
                
                chart_col4, insights_col4 = st.columns([2.5, 1.5])
                
                with chart_col4:
                    corr = filtered_df[st.session_state.numerical_cols].corr()
                    fig = px.imshow(
                        corr,
                        text_auto=True,
                        title="Feature Correlations",
                        color_continuous_scale="RdBu_r",
                        aspect="auto"
                    )
                    st.plotly_chart(fig, use_container_width=True)
                
                with insights_col4:
                    # Display insights if available
                    if st.session_state.per_chart_insights.get('correlation_heatmap'):
                        st.markdown("""
                        <div style="background-color: #fff5f5 !important; border-left: 4px solid #DC3545 !important; padding: 15px !important; border-radius: 5px !important; margin-top: 20px !important;">
                            <p style="margin: 0 !important; font-size: 14px !important; font-weight: 600 !important; color: #DC3545 !important; margin-bottom: 10px !important;">💡 Key Insights</p>
                        """, unsafe_allow_html=True)
                        st.markdown(st.session_state.per_chart_insights['correlation_heatmap'])
                        st.markdown("</div>", unsafe_allow_html=True)
                    else:
                        st.info("💡 Generate Per-Chart Insights để xem insights riêng cho chart này")
            
            # Scatter matrix for multi-variate analysis
            if len(st.session_state.numerical_cols) >= 2:
                st.markdown("#### 🎯 Multi-Variate Analysis")
                selected_cols = st.multiselect(
                    "Select columns for scatter matrix:",
                    st.session_state.numerical_cols,
                    default=st.session_state.numerical_cols[:min(3, len(st.session_state.numerical_cols))],
                    key="scatter_matrix_cols"
                )
                
                if len(selected_cols) >= 2:
                    fig = px.scatter_matrix(
                        filtered_df[selected_cols],
                        title="Scatter Matrix",
                        height=600
                    )
                    st.plotly_chart(fig, use_container_width=True)
        
        with dashboard_tabs[2]:  # KPIs & Metrics
            st.markdown("### 🎯 KPIs & Metrics Dashboard")
            st.markdown("*Dựa trên OQMLB Metrics Section*")
            
            # Show all numerical metrics
            if st.session_state.numerical_cols:
                st.markdown("#### 📊 All Metrics Summary")
                
                metrics_data = []
                for col in st.session_state.numerical_cols:
                    metrics_data.append({
                        'Metric': col.replace('_', ' ').title(),
                        'Mean': f"{df[col].mean():.2f}",
                        'Median': f"{df[col].median():.2f}",
                        'Std Dev': f"{df[col].std():.2f}",
                        'Min': f"{df[col].min():.2f}",
                        'Max': f"{df[col].max():.2f}"
                    })
                
                metrics_df = pd.DataFrame(metrics_data)
                st.dataframe(metrics_df, use_container_width=True, hide_index=True)
                
                st.markdown("---")
                
                # Individual metric deep dive
                st.markdown("#### 🔎 Metric Deep Dive")
                selected_metric = st.selectbox(
                    "Select metric to analyze:",
                    st.session_state.numerical_cols,
                    key="metric_deepdive"
                )
                
                if selected_metric:
                    col1, col2 = st.columns(2)
                    
                    with col1:
                        # Box plot
                        metric_label = selected_metric.replace('_', ' ').title()
                        fig = px.box(
                            df,
                            y=selected_metric,
                            title=f"{selected_metric} - Box Plot",
                            labels={selected_metric: metric_label}
                        )
                        st.plotly_chart(fig, use_container_width=True)
                    
                    with col2:
                        # Violin plot
                        metric_label = selected_metric.replace('_', ' ').title()
                        fig = px.violin(
                            df,
                            y=selected_metric,
                            title=f"{selected_metric} - Distribution",
                            labels={selected_metric: metric_label}
                        )
                        st.plotly_chart(fig, use_container_width=True)
        
        with dashboard_tabs[3]:  # Export & Code
            st.markdown("### 📥 Export Dashboard & Code")
            
            st.success("✅ Dashboard đã sẵn sàng để export!")
            
            # Generate Executive Summary first (if not already generated)
            client = init_gemini()
            if client and not st.session_state.executive_summary:
                if st.button("🎯 Generate Executive Summary (AI-Powered)", type="primary", use_container_width=True):
                    with st.spinner("🤖 AI đang tạo Executive Summary chuyên nghiệp..."):
                        success, summary = generate_executive_summary(
                            client, 
                            df, 
                            st.session_state.dashboard_blueprint or "", 
                            st.session_state.analysis_goal or ""
                        )
                        if success:
                            st.session_state.executive_summary = summary
                            st.rerun()
                        else:
                            st.error(f"❌ Lỗi: {summary}")
            
            # Show Executive Summary if available
            if st.session_state.executive_summary:
                with st.expander("📊 Executive Summary Preview", expanded=False):
                    st.markdown(st.session_state.executive_summary)
                st.success("✅ Executive Summary sẵn sàng cho PDF/PowerPoint export!")
            else:
                st.info("💡 Click nút 'Generate Executive Summary' để tạo báo cáo 1-page cho CEO/CFO")
            
            st.markdown("---")
            st.markdown("### 📄 Professional Reports (NEW!)")
            
            # PDF + PowerPoint Export (Row 1)
            report_col1, report_col2 = st.columns(2)
            
            with report_col1:
                st.markdown("#### 📄 PDF Report")
                st.caption("Professional multi-page report với Executive Summary, KPIs, Blueprint")
                
                if st.button("📥 Generate & Download PDF Report", use_container_width=True):
                    with st.spinner("📄 Generating PDF..."):
                        try:
                            pdf_bytes = generate_pdf_report(
                                df,
                                st.session_state.executive_summary or "Executive Summary not generated yet.",
                                st.session_state.dashboard_blueprint or "",
                                st.session_state.analysis_goal or ""
                            )
                            st.download_button(
                                "⬇️ Download PDF Report",
                                data=pdf_bytes,
                                file_name=f"analytics_report_{datetime.now().strftime('%Y%m%d_%H%M')}.pdf",
                                mime="application/pdf",
                                use_container_width=True
                            )
                            st.success("✅ PDF Report generated successfully!")
                        except Exception as e:
                            st.error(f"❌ Error generating PDF: {str(e)}")
                
                if not st.session_state.executive_summary:
                    st.warning("⚠️ Generate Executive Summary first for best results!")
            
            with report_col2:
                st.markdown("#### 📊 PowerPoint Presentation")
                st.caption("8-slide presentation ready cho Board meetings")
                
                if st.button("📥 Generate & Download PowerPoint", use_container_width=True):
                    with st.spinner("📊 Generating PowerPoint..."):
                        try:
                            ppt_bytes = generate_powerpoint_report(
                                df,
                                st.session_state.executive_summary or "Executive Summary not generated yet.",
                                st.session_state.dashboard_blueprint or "",
                                st.session_state.analysis_goal or ""
                            )
                            st.download_button(
                                "⬇️ Download PowerPoint",
                                data=ppt_bytes,
                                file_name=f"analytics_presentation_{datetime.now().strftime('%Y%m%d_%H%M')}.pptx",
                                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                                use_container_width=True
                            )
                            st.success("✅ PowerPoint generated successfully!")
                        except Exception as e:
                            st.error(f"❌ Error generating PowerPoint: {str(e)}")
                
                if not st.session_state.executive_summary:
                    st.warning("⚠️ Generate Executive Summary first for best results!")
            
            st.markdown("---")
            st.markdown("### 📦 Data & Code Exports")
            
            col1, col2, col3 = st.columns(3)
            
            with col1:
                st.markdown("#### 📊 Data Export")
                # Export cleaned data
                csv = df.to_csv(index=False).encode('utf-8')
                st.download_button(
                    "⬇️ Download Cleaned Data (CSV)",
                    data=csv,
                    file_name=f"cleaned_data_{datetime.now().strftime('%Y%m%d_%H%M')}.csv",
                    mime="text/csv",
                    use_container_width=True
                )
                
                # Export to Excel
                output = io.BytesIO()
                with pd.ExcelWriter(output, engine='xlsxwriter') as writer:
                    df.to_excel(writer, sheet_name='Data', index=False)
                excel_data = output.getvalue()
                
                st.download_button(
                    "⬇️ Download Cleaned Data (Excel)",
                    data=excel_data,
                    file_name=f"cleaned_data_{datetime.now().strftime('%Y%m%d_%H%M')}.xlsx",
                    mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                    use_container_width=True
                )
            
            with col2:
                st.markdown("#### 📋 Blueprint Export")
                if st.session_state.dashboard_blueprint:
                    st.download_button(
                        "⬇️ Download Blueprint (Markdown)",
                        data=st.session_state.dashboard_blueprint,
                        file_name=f"dashboard_blueprint_{datetime.now().strftime('%Y%m%d_%H%M')}.md",
                        mime="text/markdown",
                        use_container_width=True
                    )
                    
                    st.download_button(
                        "⬇️ Download Blueprint (Text)",
                        data=st.session_state.dashboard_blueprint,
                        file_name=f"dashboard_blueprint_{datetime.now().strftime('%Y%m%d_%H%M')}.txt",
                        mime="text/plain",
                        use_container_width=True
                    )
            
            with col3:
                st.markdown("#### 🐍 Code Export")
                
                # Generate sample dashboard code
                dashboard_code = f"""# Auto-Generated Dashboard Code
# Generated from DataAnalytics Vietnam
# Date: {datetime.now().strftime('%Y-%m-%d %H:%M')}

import streamlit as st
import pandas as pd
import plotly.express as px

# Load data
df = pd.read_csv('cleaned_data.csv')

# Dashboard Title
st.title('📊 Auto-Generated Dashboard')

# KPIs
st.markdown('## 🎯 Key Metrics')
cols = st.columns(4)

numerical_cols = {st.session_state.numerical_cols}

for i, col in enumerate(numerical_cols[:4]):
    with cols[i]:
        st.metric(
            label=col,
            value=f"{{df[col].mean():.2f}}"
        )

# Visualizations
st.markdown('## 📈 Visualizations')

# Trend chart
fig = px.line(df, y=numerical_cols[0])
st.plotly_chart(fig, use_container_width=True)

# Distribution
fig = px.histogram(df, x=numerical_cols[1] if len(numerical_cols) > 1 else numerical_cols[0])
st.plotly_chart(fig, use_container_width=True)
"""
                
                st.download_button(
                    "⬇️ Download Dashboard Code (.py)",
                    data=dashboard_code,
                    file_name=f"dashboard_code_{datetime.now().strftime('%Y%m%d_%H%M')}.py",
                    mime="text/x-python",
                    use_container_width=True
                )
                
                st.info("💡 Code này có thể chạy độc lập với Streamlit!")
            
            st.markdown("---")
            st.markdown("### 🎯 PHASE 1B: Dashboard Templates (Reusable Layouts)")
            st.caption("*Save dashboard configuration để reuse với data mới - giống Bricks.ai template system!*")
            
            template_col1, template_col2 = st.columns(2)
            
            with template_col1:
                st.markdown("#### 💾 Save Template")
                st.write("Lưu cấu hình dashboard này (blueprint + column mappings) để reuse sau")
                
                if st.session_state.blueprint_structured_data:
                    try:
                        template_json = save_dashboard_template(
                            st.session_state.blueprint_structured_data,
                            st.session_state.analysis_goal,
                            st.session_state.categorical_cols,
                            st.session_state.numerical_cols,
                            st.session_state.date_cols
                        )
                        
                        st.download_button(
                            "📥 Download Dashboard Template",
                            data=template_json,
                            file_name=f"dashboard_template_{datetime.now().strftime('%Y%m%d_%H%M')}.json",
                            mime="application/json",
                            use_container_width=True,
                            type="primary",
                            key="download_template_direct"
                        )
                        
                        st.success("✅ Click button above để download template JSON!")
                    except Exception as e:
                        st.error(f"❌ Error generating template: {str(e)}")
                else:
                    st.warning("⚠️ No blueprint data available - generate blueprint first in Step 4!")
                
                st.info("💡 Template bao gồm: Blueprint OQMLB + Column types + Analysis goals")
            
            with template_col2:
                st.markdown("#### 📂 Load Template")
                st.write("Upload template để apply lên new dataset (same structure)")
                
                uploaded_template = st.file_uploader(
                    "Upload Dashboard Template (JSON)",
                    type=['json'],
                    key="template_uploader",
                    help="Upload template JSON file đã save trước đó"
                )
                
                if uploaded_template:
                    try:
                        template_content = uploaded_template.read().decode('utf-8')
                        success, template_data = load_dashboard_template(template_content)
                        
                        if success:
                            st.success("✅ Template loaded successfully!")
                            
                            with st.expander("📋 Template Preview", expanded=False):
                                st.json({
                                    "created_at": template_data.get('created_at'),
                                    "analysis_goal": template_data.get('analysis_goal'),
                                    "categorical_cols": template_data.get('column_types', {}).get('categorical', []),
                                    "numerical_cols": template_data.get('column_types', {}).get('numerical', [])
                                })
                            
                            if st.button("✨ Apply Template to Current Data", use_container_width=True):
                                # Apply template to current dataset
                                st.session_state.dashboard_template = template_data
                                st.session_state.analysis_goal = template_data.get('analysis_goal', '')
                                st.session_state.blueprint_structured_data = template_data.get('blueprint', {})
                                st.success("🎉 Template applied! Dashboard rebuilt với template layout")
                                st.rerun()
                        else:
                            st.error(f"❌ Error loading template: {template_data.get('error', 'Unknown error')}")
                    except Exception as e:
                        st.error(f"❌ Error reading template: {str(e)}")
                
                st.info("💡 Template matching: Upload data với cùng column structure")
            
            st.markdown("---")
            st.markdown("### 🎨 Next Steps")
            
            st.markdown("""
            **Dashboard của bạn đã sẵn sàng!** 🎉
            
            Bạn có thể:
            1. ✅ **Export data** để sử dụng trong tools khác (Excel, PowerBI, Tableau)
            2. ✅ **Download blueprint** để tham khảo thiết kế chi tiết
            3. ✅ **Download Python code** để customize và deploy độc lập
            4. ✅ **Share dashboard** này với team bằng cách publish app
            
            **💡 Pro Tips:**
            - Blueprint chứa đầy đủ OQMLB framework để bạn build trên bất kỳ platform nào
            - Python code có thể customize thêm filters, charts, và features
            - Cleaned data đã được optimize và sẵn sàng cho analysis
            """)

# Footer
st.markdown("---")
st.markdown("""
<div style='text-align: center; color: #666; padding: 2rem;'>
    <p><strong>DataAnalytics Vietnam</strong> - Powered by Gemini AI</p>
    <p>Built with ❤️ for Vietnamese SMEs</p>
</div>
""", unsafe_allow_html=True)
