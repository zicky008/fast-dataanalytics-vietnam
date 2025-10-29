"""
Export utilities for DataAnalytics Vietnam
Generate professional PDF and PowerPoint reports

Chart Design Standards:
- ColorBrewer palettes (Tableau 10) - colorblind-safe, print-friendly
- WCAG AA compliance (4.5:1 text, 3:1 graphics)
- Edward Tufte principles (data-ink ratio, minimize chartjunk)
- Stephen Few dashboard design (reserved colors, high contrast)
- 300 DPI publication quality

Research: See CHART_VISUALIZATION_RESEARCH.md
"""

import io
import base64
from datetime import datetime
from typing import Dict, List, Any
import plotly.graph_objects as go
import os
import tempfile
import re

# ============================================================================
# PROFESSIONAL COLOR PALETTES (ColorBrewer - Research-Validated)
# ============================================================================

# Tableau 10 - Colorblind-safe, print-friendly (Qualitative palette)
# Source: ColorBrewer / Tableau Research
# Use for: Pie charts, categorical bar charts, distinct categories
TABLEAU_10_COLORS = [
    '#4E79A7',  # Blue (primary)
    '#F28E2B',  # Orange
    '#E15759',  # Red
    '#76B7B2',  # Teal
    '#59A14F',  # Green
    '#EDC948',  # Yellow
    '#B07AA1',  # Purple
    '#FF9DA7',  # Pink
    '#9C755F',  # Brown
    '#BAB0AC'   # Gray
]

# Sequential Blue (ColorBrewer "Blues") - for ordered data
# Use for: Heatmaps, intensity charts, ordered values
SEQUENTIAL_BLUES = ['#EFF3FF', '#C6DBEF', '#9ECAE1', '#6BAED6',
                    '#4292C6', '#2171B5', '#084594']

# Diverging Red-Blue (ColorBrewer "RdBu") - for above/below comparisons
# Use for: Benchmark comparisons, profit/loss, above/below average
DIVERGING_RDBU = ['#B2182B', '#D6604D', '#F4A582', '#FDDBC7',  # Reds
                  '#D1E5F0', '#92C5DE', '#4393C3', '#2166AC']  # Blues

# Patterns for print/colorblind accessibility (Stephen Few principle)
# Use for: Pie charts, bar charts when color alone isn't sufficient
CHART_PATTERNS = ['', '///', '...', 'xxx', '\\\\\\', '|||', '---', '+++', 'ooo', '***']

# ============================================================================


def remove_emoji(text: str) -> str:
    """
    Remove emoji characters from text for PDF compatibility.
    DejaVuSans font doesn't support emoji, causing encoding errors.

    Args:
        text: Input string that may contain emoji

    Returns:
        String with emoji removed
    """
    # Emoji pattern - matches most common emoji ranges
    emoji_pattern = re.compile("["
        u"\U0001F600-\U0001F64F"  # emoticons
        u"\U0001F300-\U0001F5FF"  # symbols & pictographs
        u"\U0001F680-\U0001F6FF"  # transport & map symbols
        u"\U0001F1E0-\U0001F1FF"  # flags (iOS)
        u"\U00002702-\U000027B0"
        u"\U000024C2-\U0001F251"
        "]+", flags=re.UNICODE)
    return emoji_pattern.sub('', text).strip()


def sanitize_text_for_pdf(text: str) -> str:
    """
    ✅ FIX #9: Sanitize text for PDF display - clean formatting issues
    Removes emoji, fixes spacing, normalizes paths to forward slashes
    
    Args:
        text: Input string that may have formatting issues
    
    Returns:
        Clean string suitable for PDF display
    """
    if not isinstance(text, str):
        text = str(text)
    
    # Remove emoji
    text = remove_emoji(text)
    
    # ✅ FIX #9: Replace backslashes with forward slashes (for paths/filenames)
    text = text.replace('\\', '/')
    
    # ✅ FIX #9: Fix spacing after special characters
    text = text.replace('%)', '%) ')  # "OEE%)83" → "OEE%) 83"
    text = text.replace('  ', ' ')  # Remove double spaces
    
    return text.strip()


def export_to_pdf(result: Dict[str, Any], df: Any, lang: str = "vi") -> bytes:
    """
    Export analysis results to professional PDF report
    
    Args:
        result: Pipeline result dictionary
        df: Original dataframe
        lang: Language code ('vi' or 'en')
    
    Returns:
        PDF file as bytes
    """
    try:
        from reportlab.lib.pagesizes import A4, letter
        from reportlab.lib.units import inch
        from reportlab.lib import colors
        from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
        from reportlab.platypus import SimpleDocTemplate, Table, TableStyle, Paragraph, Spacer, PageBreak, Image
        from reportlab.pdfgen import canvas
        from reportlab.lib.enums import TA_CENTER, TA_LEFT, TA_RIGHT
        from reportlab.pdfbase import pdfmetrics
        from reportlab.pdfbase.ttfonts import TTFont

        # Register Vietnamese-compatible font (DejaVu Sans supports Vietnamese)
        try:
            # Try system DejaVu fonts first (Linux)
            dejavu_paths = [
                '/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf',
                '/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf'
            ]

            if os.path.exists(dejavu_paths[0]):
                pdfmetrics.registerFont(TTFont('DejaVuSans', dejavu_paths[0]))
                pdfmetrics.registerFont(TTFont('DejaVuSans-Bold', dejavu_paths[1]))
                base_font = 'DejaVuSans'
                bold_font = 'DejaVuSans-Bold'
                print("✅ Vietnamese fonts loaded successfully")
            else:
                # Try without path (might be in reportlab's fonts)
                try:
                    pdfmetrics.registerFont(TTFont('DejaVuSans', 'DejaVuSans.ttf'))
                    pdfmetrics.registerFont(TTFont('DejaVuSans-Bold', 'DejaVuSans-Bold.ttf'))
                    base_font = 'DejaVuSans'
                    bold_font = 'DejaVuSans-Bold'
                    print("✅ Vietnamese fonts loaded from reportlab")
                except:
                    raise Exception("DejaVu fonts not found")
        except Exception as font_error:
            # Fallback to Helvetica if DejaVu not available (will have font issues with Vietnamese)
            print(f"⚠️ Warning: DejaVu fonts not found ({font_error}). Vietnamese characters may not display correctly.")
            print("   To fix: Install fonts with 'apt-get install fonts-dejavu' or 'yum install dejavu-sans-fonts'")
            base_font = 'Helvetica'
            bold_font = 'Helvetica-Bold'

        # Create PDF buffer
        buffer = io.BytesIO()
        doc = SimpleDocTemplate(buffer, pagesize=A4, topMargin=0.5*inch, bottomMargin=0.5*inch)

        # Get styles
        styles = getSampleStyleSheet()
        title_style = ParagraphStyle(
            'CustomTitle',
            parent=styles['Heading1'],
            fontName=bold_font,
            fontSize=24,
            textColor=colors.HexColor('#1E40AF'),
            spaceAfter=12,
            alignment=TA_CENTER
        )

        heading_style = ParagraphStyle(
            'CustomHeading',
            parent=styles['Heading2'],
            fontName=bold_font,
            fontSize=16,
            textColor=colors.HexColor('#1E40AF'),
            spaceAfter=10
        )

        normal_style = ParagraphStyle(
            'CustomNormal',
            parent=styles['Normal'],
            fontName=base_font,
            fontSize=10
        )
        
        # Build PDF content
        content = []

        # ✅ LOGO / BRANDING SUPPORT (Professional brand identity)
        # Check for logo file in project root or assets folder
        logo_paths = [
            'logo.png',
            'logo.jpg',
            'assets/logo.png',
            'assets/logo.jpg',
            'static/logo.png',
            'static/logo.jpg'
        ]

        logo_added = False
        for logo_path in logo_paths:
            if os.path.exists(logo_path):
                try:
                    # Add logo at top of report
                    logo_img = Image(logo_path, width=2.5*inch, height=0.8*inch)
                    logo_img.hAlign = 'CENTER'
                    content.append(logo_img)
                    content.append(Spacer(1, 0.2*inch))
                    logo_added = True
                    print(f"✅ Logo added from: {logo_path}")
                    break
                except Exception as logo_error:
                    print(f"⚠️ Could not load logo from {logo_path}: {str(logo_error)[:50]}")

        if not logo_added:
            # ✅ LEAN SOLUTION: Professional text-based branding (no logo image needed)
            # Create elegant header box with brand identity
            brand_style = ParagraphStyle(
                'BrandHeader',
                parent=styles['Normal'],
                fontName=bold_font,
                fontSize=10,
                textColor=colors.HexColor('#FFFFFF'),
                alignment=TA_CENTER,
                leading=14
            )

            if lang == "vi":
                brand_text = "DataAnalytics Vietnam | Phân Tích Dữ Liệu Chuyên Nghiệp"
            else:
                brand_text = "DataAnalytics Vietnam | Professional Business Intelligence"

            brand_header = Paragraph(brand_text, brand_style)

            # Create colored background box for brand header
            brand_table = Table([[brand_header]], colWidths=[7*inch])
            brand_table.setStyle(TableStyle([
                ('BACKGROUND', (0, 0), (-1, -1), colors.HexColor('#1E40AF')),
                ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
                ('VALIGN', (0, 0), (-1, -1), 'MIDDLE'),
                ('TOPPADDING', (0, 0), (-1, -1), 10),
                ('BOTTOMPADDING', (0, 0), (-1, -1), 10),
                ('LEFTPADDING', (0, 0), (-1, -1), 20),
                ('RIGHTPADDING', (0, 0), (-1, -1), 20),
            ]))

            content.append(brand_table)
            content.append(Spacer(1, 0.3*inch))
            print("✅ Professional text-based branding header added (lean solution - no logo file needed)")

        # Title (remove emoji for PDF compatibility)
        if lang == "vi":
            title = Paragraph("BÁO CÁO PHÂN TÍCH DỮ LIỆU", title_style)
            subtitle_style = ParagraphStyle('Subtitle', parent=normal_style, fontSize=9, textColor=colors.HexColor('#64748B'), alignment=TA_CENTER, fontName=base_font)
            subtitle = Paragraph("Powered by AI-Driven Analytics & ISO 8000 Standards", subtitle_style)
        else:
            title = Paragraph("DATA ANALYSIS REPORT", title_style)
            subtitle_style = ParagraphStyle('Subtitle', parent=normal_style, fontSize=9, textColor=colors.HexColor('#64748B'), alignment=TA_CENTER, fontName=base_font)
            subtitle = Paragraph("Powered by AI-Driven Analytics & ISO 8000 Standards", subtitle_style)

        content.append(title)
        content.append(subtitle)
        content.append(Spacer(1, 0.35*inch))  # ✅ Professional section separation

        # Report metadata (⭐ Enhanced with dataset profile)
        # Calculate dataset metrics
        num_rows = len(df)
        num_cols = len(df.columns)
        num_numeric = len(df.select_dtypes(include=['number']).columns)
        completeness = (1 - df.isnull().sum().sum() / (num_rows * num_cols)) * 100

        # ⭐ 5-STAR FIX: Robust multi-layer currency detection
        # Strategy: Check ALL KPIs, not just first one (Fix for Currency Inconsistency Issue #1)
        kpis_preview = result['dashboard'].get('kpis', {})
        currency_used = "USD"  # Default fallback
        
        if kpis_preview:
            # Layer 1: Check for explicit currency indicators in KPI names
            for kpi_name, kpi_data in kpis_preview.items():
                kpi_name_lower = kpi_name.lower()
                if 'vnd' in kpi_name_lower or 'vnđ' in kpi_name_lower or 'việt nam đồng' in kpi_name_lower:
                    currency_used = "VND"
                    break
                elif 'usd' in kpi_name_lower or 'dollar' in kpi_name_lower or '$' in kpi_name:
                    currency_used = "USD"
                    break
            
            # Layer 2: If no explicit indicator, use smart heuristic on ALL KPIs
            if currency_used == "USD":  # Still default, no explicit indicator found
                max_kpi_value = 0
                for kpi_name, kpi_data in kpis_preview.items():
                    kpi_value = kpi_data.get('value', 0)
                    # Focus on cost/revenue/financial KPIs for currency detection
                    if any(keyword in kpi_name.lower() for keyword in ['cost', 'revenue', 'price', 'spend', 'salary', 'income', 'expense']):
                        if abs(kpi_value) > abs(max_kpi_value):
                            max_kpi_value = kpi_value
                
                # Heuristic thresholds (refined for accuracy)
                # VND: Typically > 100K for costs, > 1M for salaries
                # USD: Typically < 10K for costs, < 100K for salaries
                if abs(max_kpi_value) > 100000:  # > 100K = likely VND
                    currency_used = "VND"
                elif 1000 < abs(max_kpi_value) <= 100000:  # 1K-100K = ambiguous, check magnitude
                    # If value is whole number > 10K, likely VND (e.g., 30000 VND vs 30.5 USD)
                    if abs(max_kpi_value) > 10000 and max_kpi_value == int(max_kpi_value):
                        currency_used = "VND"
                    else:
                        currency_used = "USD"
                else:  # < 1K = likely USD or percentage
                    currency_used = "USD"

        # Exchange rate declaration (clear and professional)
        if currency_used == "VND":
            exchange_rate = "1 USD ≈ 24,000 VND"  # ≈ symbol more professional than "approx."
        else:
            exchange_rate = "International Standard"  # More professional than "Standard: USD"

        # ✅ FIX #5: Detect and display data period (date range) for transparency
        # ISO 8000-8: Provenance - documenting data origin and temporal scope
        data_period = "Not Available"
        try:
            # Detect datetime columns
            date_cols = df.select_dtypes(include=['datetime64']).columns.tolist()
            if not date_cols:
                # Try detecting date-like string columns
                date_cols = [col for col in df.columns if any(keyword in col.lower() 
                    for keyword in ['date', 'time', 'day', 'month', 'year', 'period'])]
                if date_cols:
                    # Try converting first date column to datetime
                    import pandas as pd
                    df_temp = df.copy()
                    df_temp[date_cols[0]] = pd.to_datetime(df_temp[date_cols[0]], errors='coerce')
                    if df_temp[date_cols[0]].notna().sum() > 0:
                        date_cols = [date_cols[0]]
                        df = df_temp  # Use converted version
                    else:
                        date_cols = []
            
            # If date column found, calculate range
            if date_cols:
                date_col = date_cols[0]
                min_date = df[date_col].min()
                max_date = df[date_col].max()
                
                # Format dates (handle both datetime and string types)
                if hasattr(min_date, 'strftime'):
                    min_str = min_date.strftime("%Y-%m-%d")
                    max_str = max_date.strftime("%Y-%m-%d")
                else:
                    min_str = str(min_date)[:10]
                    max_str = str(max_date)[:10]
                
                data_period = f"{min_str} to {max_str}"
                
                # Calculate duration for additional context
                if hasattr(min_date, 'date') and hasattr(max_date, 'date'):
                    duration_days = (max_date - min_date).days
                    if duration_days > 365:
                        data_period += f" ({duration_days // 365} years, {duration_days % 365} days)"
                    elif duration_days > 30:
                        data_period += f" ({duration_days // 30} months, {duration_days % 30} days)"
                    else:
                        data_period += f" ({duration_days} days)"
        except Exception as date_error:
            # Graceful fallback if date detection fails
            print(f"⚠️ Date range detection failed: {str(date_error)[:100]}")
            data_period = "Not Available"

        # ✅ FIX #14: Wrap long text in Paragraph for Expert Perspective to prevent overflow
        expert_role_text = result['domain_info']['expert_role']
        if len(expert_role_text) > 80:
            # Use Paragraph for automatic wrapping
            expert_role_wrapped = Paragraph(expert_role_text, ParagraphStyle('ExpertWrap', 
                parent=normal_style, fontSize=9, leading=11))
        else:
            expert_role_wrapped = expert_role_text
        
        metadata_data = [
            ["Report Date" if lang == "en" else "Ngày báo cáo", datetime.now().strftime("%Y-%m-%d %H:%M:%S")],
            ["Domain" if lang == "en" else "Ngành nghề", result['domain_info']['domain_name']],
            ["Currency / Tiền tệ", f"{currency_used} ({exchange_rate})"],
            ["Data Period" if lang == "en" else "Chu kỳ dữ liệu", data_period],  # ✅ FIX #5: Added date range
            ["Expert Perspective" if lang == "en" else "Góc nhìn chuyên gia", expert_role_wrapped],  # ✅ FIX #14: Wrapped
            ["Dataset Size" if lang == "en" else "Kích thước dữ liệu", f"{num_rows:,} rows × {num_cols:,} columns ({num_numeric:,} numeric)"],
            ["Data Completeness" if lang == "en" else "Độ đầy đủ dữ liệu", f"{completeness:.1f}%"],
            ["Processing Time" if lang == "en" else "Thời gian xử lý", f"{result['performance']['total']:.1f}s"],
            ["Quality Score" if lang == "en" else "Điểm chất lượng", f"{result['quality_scores']['overall']:.0f}/100"]
        ]
        
        # ✅ FIX #14: Adjusted column widths for better text wrapping (label: 2.2", value: 4.3")
        metadata_table = Table(metadata_data, colWidths=[2.2*inch, 4.3*inch])
        metadata_table.setStyle(TableStyle([
            ('BACKGROUND', (0, 0), (0, -1), colors.HexColor('#F8FAFC')),
            ('TEXTCOLOR', (0, 0), (-1, -1), colors.HexColor('#1E293B')),
            ('ALIGN', (0, 0), (-1, -1), 'LEFT'),
            ('FONTNAME', (0, 0), (0, -1), bold_font),
            ('FONTNAME', (0, 0), (-1, -1), base_font),
            ('FONTSIZE', (0, 0), (-1, -1), 10),
            ('BOTTOMPADDING', (0, 0), (-1, -1), 8),
            ('GRID', (0, 0), (-1, -1), 0.5, colors.HexColor('#E2E8F0')),
            # ✅ FIX #14: Enable top-align and wrap for Expert Perspective row (row 4)
            ('VALIGN', (0, 4), (-1, 4), 'TOP')
        ]))

        content.append(metadata_table)
        content.append(Spacer(1, 0.4*inch))  # ✅ FIX #17: Consistent spacing

        # ✅ FIX #19: Add icons to section headers for 5-star professional experience
        # Executive Summary with icon
        exec_icon = "📊" if lang == "en" else "📊"
        exec_title = f"{exec_icon} Executive Summary" if lang == "en" else f"{exec_icon} Tóm Tắt Điều Hành"
        content.append(Paragraph(exec_title, heading_style))
        
        summary_text = result['insights'].get('executive_summary', 'No summary available')
        # ✅ FIX #9: Sanitize text for proper formatting
        summary_text = sanitize_text_for_pdf(summary_text)
        
        # ✅ FIX #19: Add highlight box for executive summary (5-star UX)
        from reportlab.platypus import KeepTogether
        summary_para = Paragraph(summary_text, ParagraphStyle('SummaryHighlight',
            parent=normal_style,
            backColor=colors.HexColor('#FEF3C7'),  # Light yellow highlight
            borderColor=colors.HexColor('#F59E0B'),  # Amber border
            borderWidth=1,
            borderPadding=12,
            spaceBefore=6,
            spaceAfter=6,
            leading=14
        ))
        content.append(summary_para)
        content.append(Spacer(1, 0.3*inch))  # ✅ FIX #17: Consistent spacing

        # Key KPIs with icon
        kpi_icon = "🎯" if lang == "en" else "🎯"
        kpi_title = f"{kpi_icon} Key Performance Indicators" if lang == "en" else f"{kpi_icon} Chỉ Số Hiệu Suất Chính"
        content.append(Paragraph(kpi_title, heading_style))

        kpis = result['dashboard'].get('kpis', {})
        if kpis:
            # ⭐ Updated to include benchmark source (addresses real user feedback)
            kpi_data = [[
                "KPI",
                "Value" if lang == "en" else "Giá trị",
                "Status" if lang == "en" else "Trạng thái",
                "Benchmark" if lang == "en" else "Chuẩn",
                "Source" if lang == "en" else "Nguồn"
            ]]

            for kpi_name, kpi_info in list(kpis.items())[:10]:  # Top 10 KPIs
                # ✅ FIX #6 & #9: Clean KPI name - remove markdown and fix spacing
                # CRITICAL: Replace longest patterns FIRST (### before ##)
                clean_kpi_name = kpi_name.replace('###', '').replace('##', '').replace('**', '').replace('#', '').strip()
                # ✅ FIX #9: Ensure proper spacing after percentage symbols
                clean_kpi_name = clean_kpi_name.replace('%)', '%) ').replace('  ', ' ').strip()
                
                # ✅ FIX #14 & #20: Wrap long KPI names in Paragraph for Vietnamese text optimization
                if len(clean_kpi_name) > 50:
                    clean_kpi_name_cell = Paragraph(clean_kpi_name, ParagraphStyle('KPIWrap', 
                        parent=normal_style, fontSize=9, leading=11, wordWrap='CJK'))
                else:
                    clean_kpi_name_cell = clean_kpi_name
                
                # ✅ FIX #18: Add clickable links to benchmark sources for depth
                source = kpi_info.get('benchmark_source', 'Industry Standard')
                
                # Map sources to URLs for interactivity
                source_urls = {
                    'McKinsey Manufacturing Report': 'https://www.mckinsey.com/capabilities/operations/our-insights',
                    'Gartner IT Benchmarks': 'https://www.gartner.com/en/research/benchmarking',
                    'WordStream PPC Benchmarks': 'https://www.wordstream.com/blog/ws/2019/11/12/google-ads-benchmarks',
                    'HubSpot Marketing Benchmarks': 'https://www.hubspot.com/marketing-statistics',
                    'Salesforce Sales Benchmarks': 'https://www.salesforce.com/resources/research-reports/',
                    'Zendesk Support Benchmarks': 'https://www.zendesk.com/benchmark/',
                    'Deloitte Financial Services': 'https://www2.deloitte.com/us/en/pages/financial-services/topics/center-for-financial-services.html',
                    'PwC HR Metrics': 'https://www.pwc.com/gx/en/services/people-organisation.html'
                }
                
                # Check if source has a URL
                source_url = None
                for known_source, url in source_urls.items():
                    if known_source.lower() in source.lower():
                        source_url = url
                        break
                
                # Create clickable link if URL exists
                if source_url:
                    source_link = f'<a href="{source_url}" color="blue"><u>{source}</u></a>'
                    source_paragraph = Paragraph(source_link, normal_style)
                else:
                    # No link, just wrap for word wrap
                    source_paragraph = Paragraph(source, normal_style) if len(source) > 25 else source

                # Format value with thousand separators
                value = kpi_info['value']
                if value >= 1000:
                    formatted_value = f"{value:,.1f}"
                else:
                    formatted_value = f"{value:.2f}"

                # Add unit if it's a salary KPI
                if 'salary' in kpi_name.lower() or 'compensation' in kpi_name.lower():
                    formatted_value += f" {currency_used}/year"

                # ✅ FIX #6: Enhanced benchmark formatting with context
                benchmark = kpi_info.get('benchmark', 'N/A')
                if benchmark != 'N/A' and isinstance(benchmark, (int, float)):
                    # Add target indicator for clarity
                    status = kpi_info.get('status', 'N/A')
                    target_symbol = ""
                    
                    # Determine if higher or lower is better based on KPI context
                    kpi_lower = clean_kpi_name.lower()
                    if any(keyword in kpi_lower for keyword in ['cost', 'expense', 'defect', 'downtime', 'error', 'reject']):
                        # Lower is better
                        target_symbol = "≤" if status in ['Below', 'Good'] else ">"
                    elif any(keyword in kpi_lower for keyword in ['revenue', 'profit', 'efficiency', 'yield', 'quality', 'conversion', 'roi', 'roas']):
                        # Higher is better
                        target_symbol = "≥" if status in ['Above', 'Good'] else "<"
                    
                    # Format number with thousands separator
                    if benchmark >= 1000:
                        formatted_benchmark = f"{target_symbol} {benchmark:,.0f}".strip()
                    else:
                        formatted_benchmark = f"{target_symbol} {benchmark:.1f}".strip()
                    
                    # Clean up if no symbol determined
                    formatted_benchmark = formatted_benchmark.strip()
                else:
                    formatted_benchmark = str(benchmark)

                # Note: source_paragraph already created above with clickable links
                
                # ✅ FIX #15 & #16: Enhanced status with CONSISTENT logic, color, and icons
                status_raw = kpi_info.get('status', 'N/A')
                status_display = status_raw
                status_color = colors.HexColor('#64748B')  # Default gray
                
                if status_raw not in ['N/A', 'Unknown']:
                    # Determine if current status is good or bad based on KPI type
                    kpi_lower = clean_kpi_name.lower() if isinstance(clean_kpi_name, str) else str(clean_kpi_name).lower()
                    is_cost_type = any(keyword in kpi_lower for keyword in ['cost', 'expense', 'defect', 'downtime', 'error', 'reject', 'waste', 'rate'])
                    is_revenue_type = any(keyword in kpi_lower for keyword in ['revenue', 'profit', 'efficiency', 'yield', 'quality', 'conversion', 'roi', 'roas', 'satisfaction', 'oee'])
                    
                    # ✅ FIX #15: CORRECTED LOGIC - For costs/defects, "Above" = bad (red), "Below" = good (green)
                    if status_raw in ['Above', 'High', 'Over']:
                        if is_cost_type:
                            status_display = f"Above ⬇️"  # Above cost/defect = BAD
                            status_color = colors.HexColor('#DC2626')  # Red (bad)
                        elif is_revenue_type:
                            status_display = f"Above ⬆️"  # Above revenue = GOOD
                            status_color = colors.HexColor('#16A34A')  # Green (good)
                        else:
                            status_display = status_raw
                    elif status_raw in ['Below', 'Low', 'Under']:
                        if is_cost_type:
                            status_display = f"Below ⬆️"  # Below cost/defect = GOOD
                            status_color = colors.HexColor('#16A34A')  # Green (good)
                        elif is_revenue_type:
                            status_display = f"Below ⬇️"  # Below revenue = BAD
                            status_color = colors.HexColor('#DC2626')  # Red (bad)
                        else:
                            status_display = status_raw
                    elif status_raw in ['Good', 'Excellent', 'On Target']:
                        status_display = f"{status_raw} ✓"
                        status_color = colors.HexColor('#16A34A')  # Green
                    elif status_raw in ['Poor', 'Critical', 'Alert']:
                        status_display = f"{status_raw} ✗"
                        status_color = colors.HexColor('#DC2626')  # Red
                
                # ✅ FIX #16: Wrap status in colored Paragraph for accessibility
                status_cell = Paragraph(f'<font color="{status_color}">{status_display}</font>', 
                    ParagraphStyle('Status', parent=normal_style, fontSize=9, alignment=1))
                
                kpi_data.append([
                    clean_kpi_name_cell,  # ✅ FIX #14: Wrapped KPI name
                    formatted_value,
                    status_cell,  # ✅ FIX #15 & #16: Colored status with consistent logic
                    formatted_benchmark,  # ✅ FIX #6: Enhanced with target indicators
                    source_paragraph  # ✅ Full source with word wrap
                ])

            # ✅ FIX #14: Further optimized column widths to prevent overflow
            # KPI: 2.0" (increased for long names), Value: 0.8", Status: 0.9", Benchmark: 0.8", Source: 2.0" (balanced)
            # Total: 6.5" (fits standard page width with margins)
            kpi_table = Table(kpi_data, colWidths=[2.0*inch, 0.8*inch, 0.9*inch, 0.8*inch, 2.0*inch])
            kpi_table.setStyle(TableStyle([
                ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor('#1E40AF')),
                ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
                ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
                ('FONTNAME', (0, 0), (-1, 0), bold_font),
                ('FONTNAME', (0, 1), (-1, -1), base_font),
                ('FONTSIZE', (0, 0), (-1, -1), 9),
                ('BOTTOMPADDING', (0, 0), (-1, -1), 8),
                ('BACKGROUND', (0, 1), (-1, -1), colors.beige),
                ('GRID', (0, 0), (-1, -1), 0.5, colors.HexColor('#E2E8F0')),
                ('ROWBACKGROUNDS', (0, 1), (-1, -1), [colors.white, colors.HexColor('#F8FAFC')]),
                # ✅ FIX #14: Enable word wrap and top alignment for KPI names (column 0) and Source (column 4)
                ('VALIGN', (0, 1), (0, -1), 'TOP'),  # Top-align KPI names
                ('ALIGN', (0, 1), (0, -1), 'LEFT'),  # Left-align KPI names for readability
                ('VALIGN', (4, 0), (4, -1), 'TOP'),  # Top-align sources
                ('LEFTPADDING', (0, 1), (0, -1), 6),  # Padding for KPI names
                ('LEFTPADDING', (4, 1), (4, -1), 6),  # Padding for sources
                ('RIGHTPADDING', (4, 1), (4, -1), 6),
            ]))
            
            content.append(kpi_table)

            # ✅ FIX #8: Show calculation breakdowns for complex/composite KPIs (reproducibility)
            # ISO 8000 Principle: Calculations must be transparent and reproducible
            complex_kpis = {}
            for kpi_name, kpi_info in kpis.items():
                if 'components' in kpi_info:
                    complex_kpis[kpi_name] = kpi_info
            
            if complex_kpis:
                content.append(Spacer(1, 0.2*inch))
                
                # Add calculation breakdown section
                calc_title = Paragraph(
                    "<b>Calculation Breakdowns</b>" if lang == "en" else "<b>Chi Tiết Tính Toán</b>",
                    ParagraphStyle('CalcTitle', parent=normal_style, fontSize=11, fontName=bold_font, textColor=colors.HexColor('#1E40AF'))
                )
                content.append(calc_title)
                content.append(Spacer(1, 0.1*inch))
                
                # Display each complex KPI's formula and components
                for kpi_name, kpi_info in list(complex_kpis.items())[:3]:  # Top 3 complex KPIs
                    clean_name = kpi_name.replace('##', '').replace('###', '').replace('**', '').strip()
                    components = kpi_info['components']
                    
                    # Build formula display
                    if 'OEE' in kpi_name or 'Equipment Effectiveness' in kpi_name:
                        formula = "OEE = Availability × Performance × Quality"
                        component_lines = [
                            f"  • Availability: {components.get('Availability', 0):.1f}%",
                            f"  • Performance: {components.get('Performance', 0):.1f}%",
                            f"  • Quality: {components.get('Quality', 0):.1f}%"
                        ]
                    elif 'ROI' in kpi_name:
                        formula = "ROI = (Net Profit / Cost of Investment) × 100"
                        component_lines = [
                            f"  • Net Profit: {components.get('Net Profit', 0):,.0f}",
                            f"  • Investment Cost: {components.get('Investment Cost', 0):,.0f}"
                        ]
                    elif 'ROAS' in kpi_name:
                        formula = "ROAS = Revenue from Ads / Ad Spend"
                        component_lines = [
                            f"  • Revenue: {components.get('Revenue', 0):,.0f}",
                            f"  • Ad Spend: {components.get('Ad Spend', 0):,.0f}"
                        ]
                    else:
                        # Generic display
                        formula = f"{clean_name} = f(components)"
                        component_lines = [f"  • {k}: {v:.2f}" for k, v in components.items()]
                    
                    # Assemble the breakdown text
                    breakdown_text = f"<b>{clean_name}</b>: {kpi_info['value']:.2f}<br/>"
                    breakdown_text += f"<i>{formula}</i><br/>"
                    breakdown_text += "<br/>".join(component_lines)
                    
                    calc_para = Paragraph(breakdown_text, normal_style)
                    content.append(calc_para)
                    content.append(Spacer(1, 0.15*inch))

            # ⭐ Add KPI Status legend for clarity (addresses real user feedback)
            content.append(Spacer(1, 0.15*inch))
            if lang == "en":
                status_note = "<i><font size=8>Status Guide: ✅ Above = +10% vs benchmark | ➡️ Competitive = ±10% | ⚠️ Below = -10% vs benchmark. Note: Lower is better for costs/time.</font></i>"
            else:
                status_note = "<i><font size=8>Hướng dẫn Trạng thái: ✅ Trên chuẩn = +10% so với benchmark | ➡️ Cạnh tranh = ±10% | ⚠️ Dưới chuẩn = -10% so với benchmark. Lưu ý: Thấp hơn tốt hơn cho chi phí/thời gian.</font></i>"
            content.append(Paragraph(status_note, normal_style))

        content.append(Spacer(1, 0.5*inch))  # ✅ Clear section transition

        # ✅ FIX #19: Key Insights with icon
        insights_icon = "💡" if lang == "en" else "💡"
        insights_title = f"{insights_icon} Key Insights" if lang == "en" else f"{insights_icon} Insights Chính"
        content.append(Paragraph(insights_title, heading_style))
        content.append(Spacer(1, 0.2*inch))  # ✅ FIX #17: Consistent spacing

        for i, insight in enumerate(result['insights'].get('key_insights', [])[:5], 1):
            # Use text labels instead of emoji for PDF compatibility
            impact_label = "[HIGH]" if insight['impact'] == 'high' else "[MEDIUM]" if insight['impact'] == 'medium' else "[LOW]"
            # ✅ FIX #9: Sanitize insight text
            title_clean = sanitize_text_for_pdf(insight['title'])
            desc_clean = sanitize_text_for_pdf(insight['description'])
            insight_text = f"{impact_label} <b>{title_clean}</b><br/>{desc_clean}"
            content.append(Paragraph(insight_text, normal_style))
            content.append(Spacer(1, 0.15*inch))  # ✅ FIX #17: Consistent spacing

        content.append(Spacer(1, 0.3*inch))  # ✅ FIX #17: Consistent section break

        # ✅ FIX #19: Recommendations with icon
        rec_icon = "✨" if lang == "en" else "✨"
        rec_title = f"{rec_icon} Recommendations" if lang == "en" else f"{rec_icon} Khuyến Nghị"
        content.append(Paragraph(rec_title, heading_style))
        content.append(Spacer(1, 0.2*inch))  # ✅ FIX #17: Consistent spacing

        for i, rec in enumerate(result['insights'].get('recommendations', [])[:5], 1):
            # Use text labels instead of emoji for PDF compatibility
            priority_label = "[HIGH]" if rec['priority'] == 'high' else "[MEDIUM]" if rec['priority'] == 'medium' else "[LOW]"
            
            # ✅ FIX #11: Determine responsible role based on action keywords and domain
            responsible_role = rec.get('responsible', None)  # Check if already provided
            if not responsible_role:
                # Infer from action content and domain
                action_lower = rec['action'].lower()
                domain_lower = result['domain_info']['domain'].lower()
                
                if any(keyword in action_lower for keyword in ['budget', 'cost', 'financial', 'investment', 'roi', 'pricing']):
                    responsible_role = "CFO / Finance Director"
                elif any(keyword in action_lower for keyword in ['marketing', 'campaign', 'advertising', 'social media', 'brand', 'customer acquisition']):
                    responsible_role = "CMO / Marketing Manager"
                elif any(keyword in action_lower for keyword in ['hire', 'training', 'employee', 'hr', 'talent', 'compensation', 'culture']):
                    responsible_role = "CHRO / HR Manager"
                elif any(keyword in action_lower for keyword in ['production', 'manufacturing', 'quality', 'oee', 'process', 'equipment', 'downtime']):
                    responsible_role = "COO / Operations Manager"
                elif any(keyword in action_lower for keyword in ['technology', 'system', 'software', 'automation', 'digital', 'integration']):
                    responsible_role = "CTO / IT Manager"
                elif any(keyword in action_lower for keyword in ['sales', 'pipeline', 'lead', 'conversion', 'deal', 'revenue']):
                    responsible_role = "VP Sales / Sales Manager"
                elif any(keyword in action_lower for keyword in ['customer service', 'support', 'satisfaction', 'resolution', 'ticket']):
                    responsible_role = "Customer Service Manager"
                elif any(keyword in action_lower for keyword in ['data', 'analytics', 'dashboard', 'metric', 'tracking', 'reporting']):
                    responsible_role = "Data Analyst / Business Intelligence"
                else:
                    # Default based on domain
                    if 'manufacturing' in domain_lower:
                        responsible_role = "Operations Manager"
                    elif 'marketing' in domain_lower or 'advertising' in domain_lower:
                        responsible_role = "Marketing Manager"
                    elif 'hr' in domain_lower or 'human' in domain_lower:
                        responsible_role = "HR Manager"
                    elif 'finance' in domain_lower:
                        responsible_role = "Finance Manager"
                    elif 'sales' in domain_lower:
                        responsible_role = "Sales Manager"
                    else:
                        responsible_role = "Department Head"
            
            rec_text = f"{priority_label} <b>{rec['action']}</b><br/>Expected Impact: {rec['expected_impact']}<br/>Timeline: {rec['timeline']}<br/><b>Responsible:</b> {responsible_role}"
            content.append(Paragraph(rec_text, normal_style))
            content.append(Spacer(1, 0.18*inch))  # ✅ Tight, professional spacing

        # ✅ FIX #17: Consistent spacing before page break
        content.append(Spacer(1, 0.4*inch))  # Consistent breathing room
        content.append(PageBreak())

        # ✅ FIX #19: Charts with icon
        chart_icon = "📈" if lang == "en" else "📈"
        chart_title = f"{chart_icon} Visual Analysis" if lang == "en" else f"{chart_icon} Phân Tích Trực Quan"
        content.append(Paragraph(chart_title, heading_style))
        content.append(Spacer(1, 0.2*inch))  # ✅ FIX #17: Consistent spacing

        charts = result['dashboard']['charts']
        charts_exported = 0
        total_charts = len(charts)

        print(f"\n{'='*80}")
        print(f"📊 CHART EXPORT DEBUG INFO")
        print(f"{'='*80}")
        print(f"Total charts in dashboard: {total_charts}")
        print(f"Starting export process...")
        print(f"{'='*80}\n")

        for i, chart in enumerate(charts):  # Export ALL charts
            print(f"\n--- Processing Chart {i+1}/{total_charts} ---")
            try:
                # Add chart title
                chart_title = chart.get('title', f'Chart {i+1}')
                chart_type = chart.get('figure', {}).data[0].type if chart.get('figure') and chart.get('figure').data else 'unknown'
                print(f"  Title: {chart_title}")
                print(f"  Type: {chart_type}")

                # ✅ Professional chart title with clear visual hierarchy
                chart_title_style = ParagraphStyle(
                    'ChartTitle',
                    parent=normal_style,
                    fontSize=12,
                    textColor=colors.HexColor('#1E40AF'),
                    spaceAfter=6,
                    fontName=bold_font
                )
                title_para = Paragraph(f"<b>{chart_title}</b>", chart_title_style)
                content.append(title_para)
                content.append(Spacer(1, 0.12*inch))  # ✅ Tight title-chart connection

                # Convert Plotly chart to image using kaleido
                fig = chart['figure']

                # Try to export chart image - Multiple fallback methods
                chart_exported = False

                # Method 1: Try kaleido (requires Chrome/Chromium)
                if not chart_exported:
                    try:
                        img_bytes = fig.to_image(format="png", width=800, height=450, engine="kaleido")
                        img_stream = io.BytesIO(img_bytes)
                        img = Image(img_stream, width=6.5*inch, height=3.6*inch)
                        content.append(img)
                        charts_exported += 1
                        chart_exported = True
                        print(f"✅ Successfully exported chart {i+1}: {chart_title}")
                    except Exception as img_error:
                        print(f"⚠️ Chart {i+1} kaleido method failed: {str(img_error)[:80]}")

                # Method 2: Try write_image to temp file
                if not chart_exported:
                    try:
                        with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp:
                            fig.write_image(tmp.name, format='png', width=800, height=450)
                            img = Image(tmp.name, width=6.5*inch, height=3.6*inch)
                            content.append(img)
                            charts_exported += 1
                            chart_exported = True
                            os.unlink(tmp.name)
                            print(f"✅ Successfully exported chart {i+1} (write_image): {chart_title}")
                    except Exception as alt_error:
                        print(f"⚠️ Chart {i+1} write_image method failed: {str(alt_error)[:80]}")

                # Method 2.5: Matplotlib fallback (NO Chrome/Kaleido needed!)
                # ⚠️ Enhanced for HIGHER QUALITY (300 DPI, better styling)
                if not chart_exported:
                    try:
                        import matplotlib
                        matplotlib.use('Agg')  # Non-interactive backend
                        import matplotlib.pyplot as plt
                        from io import BytesIO

                        # ✅ CREATE PROFESSIONAL MATPLOTLIB FIGURE
                        # 300 DPI = publication quality (ColorBrewer/academic standard)
                        mpl_fig = plt.figure(figsize=(10, 5.6), dpi=300)
                        ax = mpl_fig.add_subplot(111)

                        # ✅ PROFESSIONAL STYLING (Tufte/Few principles)
                        # Clean, minimal style - let data shine
                        plt.style.use('seaborn-v0_8-whitegrid')  # Clean background
                        mpl_fig.patch.set_facecolor('#FFFFFF')  # Pure white for print
                        ax.set_facecolor('#FAFAFA')  # Slight gray to separate from page

                        # ✅ REMOVE CHARTJUNK (Edward Tufte principle)
                        # Remove unnecessary visual elements that don't encode data
                        ax.spines['top'].set_visible(False)    # Remove top border
                        ax.spines['right'].set_visible(False)  # Remove right border
                        ax.spines['left'].set_color('#CCCCCC')   # Subtle left axis
                        ax.spines['bottom'].set_color('#CCCCCC') # Subtle bottom axis
                        ax.spines['left'].set_linewidth(0.8)
                        ax.spines['bottom'].set_linewidth(0.8)

                        # Extract data from plotly figure
                        # Enhanced conversion - supports more chart types
                        for trace_idx, trace in enumerate(fig.data):
                            trace_type = trace.type if hasattr(trace, 'type') else 'scatter'
                            trace_name = trace.name if hasattr(trace, 'name') else None

                            # Scatter/Line charts
                            if hasattr(trace, 'x') and hasattr(trace, 'y'):
                                if trace_type in ['scatter', 'scattergl']:
                                    mode = trace.mode if hasattr(trace, 'mode') else 'lines+markers'
                                    if 'lines' in mode:
                                        ax.plot(trace.x, trace.y, label=trace_name, marker='o' if 'markers' in mode else None)
                                    elif 'markers' in mode:
                                        ax.scatter(trace.x, trace.y, label=trace_name, alpha=0.6)

                                        # Add trendline for scatter plots
                                        try:
                                            import numpy as np
                                            x_numeric = [float(xi) for xi in trace.x if xi is not None]
                                            y_numeric = [float(yi) for yi in trace.y if yi is not None]

                                            if len(x_numeric) >= 2 and len(y_numeric) >= 2:
                                                # Calculate linear regression
                                                z = np.polyfit(x_numeric, y_numeric, 1)
                                                p = np.poly1d(z)
                                                ax.plot(x_numeric, p(x_numeric), "--", alpha=0.8, linewidth=2,
                                                       label=f'Trendline (y={z[0]:.2f}x+{z[1]:.2f})')

                                                # Calculate R²
                                                y_pred = p(x_numeric)
                                                ss_res = np.sum((y_numeric - y_pred) ** 2)
                                                ss_tot = np.sum((y_numeric - np.mean(y_numeric)) ** 2)
                                                r_squared = 1 - (ss_res / ss_tot) if ss_tot > 0 else 0

                                                # Add R² annotation
                                                ax.text(0.05, 0.95, f'R² = {r_squared:.3f}',
                                                       transform=ax.transAxes, fontsize=10,
                                                       verticalalignment='top',
                                                       bbox=dict(boxstyle='round', facecolor='wheat', alpha=0.5))
                                        except (ValueError, TypeError, ImportError):
                                            pass  # Skip trendline if data not suitable
                                    else:
                                        ax.plot(trace.x, trace.y, label=trace_name)

                                # Bar charts
                                elif trace_type == 'bar':
                                    # ✅ PROFESSIONAL BAR CHART (ColorBrewer + High Contrast)
                                    # Use Tableau 10 colors for multiple series
                                    bar_color = TABLEAU_10_COLORS[trace_idx % len(TABLEAU_10_COLORS)]

                                    # Handle grouped/stacked bars
                                    if trace_idx == 0:
                                        bars = ax.bar(trace.x, trace.y, label=trace_name,
                                                     color=bar_color, edgecolor='#333333', linewidth=0.8)
                                        # ✅ CLARITY: Add value labels on top of bars
                                        for bar in bars:
                                            height = bar.get_height()
                                            if height > 0:  # Only show positive values
                                                ax.text(bar.get_x() + bar.get_width()/2., height,
                                                       f'{height:.0f}',
                                                       ha='center', va='bottom', fontsize=9,
                                                       fontweight='bold', color='#000000')
                                    else:
                                        ax.bar(trace.x, trace.y, label=trace_name,
                                              color=bar_color, edgecolor='#333333', linewidth=0.8)

                                    # ✅ READABILITY: Rotate x-axis labels if they're long
                                    if hasattr(trace, 'x') and len(trace.x) > 0:
                                        if isinstance(trace.x[0], str) and len(str(trace.x[0])) > 10:
                                            plt.setp(ax.get_xticklabels(), rotation=45, ha='right', fontsize=9)

                                # Line charts
                                elif trace_type == 'line':
                                    ax.plot(trace.x, trace.y, label=trace_name)

                                # Histogram
                                elif trace_type == 'histogram':
                                    if hasattr(trace, 'x'):
                                        ax.hist(trace.x, bins=20, label=trace_name, alpha=0.7)
                                    elif hasattr(trace, 'y'):
                                        ax.hist(trace.y, bins=20, label=trace_name, alpha=0.7, orientation='horizontal')

                                # Box plot
                                elif trace_type == 'box':
                                    if hasattr(trace, 'y'):
                                        ax.boxplot([trace.y], labels=[trace_name if trace_name else ''])

                            # Pie chart (special case - no x/y)
                            elif trace_type == 'pie':
                                if hasattr(trace, 'labels') and hasattr(trace, 'values'):
                                    # ✅ PROFESSIONAL PIE CHART (ColorBrewer + Accessibility)
                                    # Extract colors from Plotly trace if available
                                    colors_list = None
                                    if hasattr(trace, 'marker') and hasattr(trace.marker, 'colors'):
                                        colors_list = trace.marker.colors

                                    # Use Tableau 10 colorblind-safe palette if no colors specified
                                    if colors_list is None:
                                        colors_list = TABLEAU_10_COLORS

                                    # Create pie chart with percentage labels
                                    wedges, texts, autotexts = ax.pie(
                                        trace.values,
                                        labels=trace.labels,
                                        autopct='%1.1f%%',
                                        colors=colors_list,
                                        startangle=90,
                                        textprops={'fontsize': 10, 'fontweight': 'bold'}
                                    )

                                    # ✅ ACCESSIBILITY: Add patterns for print/colorblind users (Stephen Few)
                                    # Patterns allow differentiation even in grayscale
                                    for i, wedge in enumerate(wedges):
                                        wedge.set_hatch(CHART_PATTERNS[i % len(CHART_PATTERNS)])
                                        wedge.set_edgecolor('#FFFFFF')  # White borders for separation
                                        wedge.set_linewidth(1.5)

                                    # ✅ HIGH CONTRAST: Make percentage text readable
                                    for autotext in autotexts:
                                        autotext.set_color('#000000')  # Black text
                                        autotext.set_fontsize(9)
                                        autotext.set_fontweight('bold')

                                    ax.axis('equal')  # Equal aspect ratio ensures circular pie

                            # Heatmap (special case - 2D data)
                            elif trace_type == 'heatmap':
                                if hasattr(trace, 'z'):
                                    im = ax.imshow(trace.z, cmap='viridis', aspect='auto')
                                    if hasattr(trace, 'x'):
                                        ax.set_xticks(range(len(trace.x)))
                                        ax.set_xticklabels(trace.x, rotation=45)
                                    if hasattr(trace, 'y'):
                                        ax.set_yticks(range(len(trace.y)))
                                        ax.set_yticklabels(trace.y)
                                    plt.colorbar(im, ax=ax)

                        # ✅ PROFESSIONAL TYPOGRAPHY (WCAG AA compliant)
                        # High contrast text for readability
                        if fig.layout.title.text:
                            ax.set_title(fig.layout.title.text,
                                       fontsize=14, fontweight='bold', pad=20,
                                       color='#000000')  # Black for 4.5:1 contrast
                        if fig.layout.xaxis.title.text:
                            ax.set_xlabel(fig.layout.xaxis.title.text,
                                        fontsize=11, fontweight='semibold',
                                        color='#333333', labelpad=10)
                        if fig.layout.yaxis.title.text:
                            ax.set_ylabel(fig.layout.yaxis.title.text,
                                        fontsize=11, fontweight='semibold',
                                        color='#333333', labelpad=10)

                        # ✅ SUBTLE GRID (Tufte: Grid should be "just noticeable")
                        # Horizontal grid only for easier reading (Stephen Few)
                        ax.grid(True, axis='y', alpha=0.3, linestyle='--',
                               linewidth=0.5, color='#CCCCCC')
                        ax.set_axisbelow(True)  # Grid behind data

                        # ✅ CLEAN TICK MARKS (remove tick lines, keep labels)
                        ax.tick_params(axis='both', which='both', length=0,  # No tick marks
                                     labelsize=9, colors='#333333')  # Readable labels
                        ax.tick_params(axis='x', pad=8)  # White space
                        ax.tick_params(axis='y', pad=8)

                        # ✅ PROFESSIONAL LEGEND (if needed)
                        if any(hasattr(trace, 'name') and trace.name for trace in fig.data):
                            ax.legend(loc='best',
                                    frameon=True,
                                    framealpha=0.95,
                                    edgecolor='#CCCCCC',
                                    fancybox=False,  # No rounded corners (cleaner)
                                    shadow=False,    # No shadow (less chartjunk)
                                    fontsize=9,
                                    labelspacing=1.0,
                                    borderpad=0.8)

                        # Tight layout for better spacing
                        plt.tight_layout()

                        # Save to bytes with HIGH DPI (300 = publication quality)
                        buf = BytesIO()
                        mpl_fig.savefig(buf, format='png', dpi=300, bbox_inches='tight',
                                       facecolor='white', edgecolor='none')
                        buf.seek(0)

                        # Add to PDF
                        img = Image(buf, width=6.5*inch, height=3.6*inch)
                        content.append(img)
                        charts_exported += 1
                        chart_exported = True
                        plt.close(mpl_fig)
                        print(f"✅ Successfully exported chart {i+1} (matplotlib fallback): {chart_title}")

                    except ImportError:
                        print(f"⚠️ Chart {i+1} matplotlib not available")
                    except Exception as mpl_error:
                        print(f"⚠️ Chart {i+1} matplotlib method failed: {str(mpl_error)[:80]}")

                # ✅ FIX #10: Method 4 - Universal text-based representation (ALWAYS works)
                # Instead of showing error message, show data in table format
                if not chart_exported:
                    try:
                        # Extract key data points from chart for text representation
                        chart_data_summary = []
                        
                        for trace in fig.data[:5]:  # Top 5 traces
                            trace_type = trace.type if hasattr(trace, 'type') else 'unknown'
                            trace_name = trace.name if hasattr(trace, 'name') else 'Data'
                            
                            if trace_type == 'pie' and hasattr(trace, 'labels') and hasattr(trace, 'values'):
                                # Pie chart - show labels and values
                                for label, value in zip(trace.labels[:5], trace.values[:5]):
                                    chart_data_summary.append([str(label), f"{value:.2f}"])
                            
                            elif hasattr(trace, 'x') and hasattr(trace, 'y'):
                                # Line/Bar/Scatter - show key points
                                x_vals = list(trace.x)[:5]  # First 5 points
                                y_vals = list(trace.y)[:5]
                                for x, y in zip(x_vals, y_vals):
                                    chart_data_summary.append([str(x), f"{y:.2f}"])
                        
                        if chart_data_summary:
                            # Create compact data table
                            summary_header = [[f"<b>{chart_title}</b> (Data Summary)", ""]]
                            data_table = Table(summary_header + chart_data_summary, colWidths=[3.5*inch, 2*inch])
                            data_table.setStyle(TableStyle([
                                ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor('#1E40AF')),
                                ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
                                ('FONTNAME', (0, 0), (-1, 0), bold_font),
                                ('FONTSIZE', (0, 0), (-1, -1), 9),
                                ('GRID', (0, 0), (-1, -1), 0.5, colors.HexColor('#E2E8F0')),
                                ('ROWBACKGROUNDS', (0, 1), (-1, -1), [colors.white, colors.HexColor('#F8FAFC')])
                            ]))
                            content.append(data_table)
                            charts_exported += 1
                            print(f"✅ Chart {i+1} exported as data table: {chart_title}")
                        else:
                            # Last resort - show message
                            if lang == "vi":
                                degraded_msg = f"<i>Biểu đồ: {chart_title} (Dữ liệu có sẵn trong phần KPIs)</i>"
                            else:
                                degraded_msg = f"<i>Chart: {chart_title} (Data available in KPIs section)</i>"
                            
                            degraded_para = Paragraph(degraded_msg, normal_style)
                            content.append(degraded_para)
                            print(f"ℹ️ Chart {i+1} noted without visual: {chart_title}")
                    
                    except Exception as fallback_error:
                        print(f"⚠️ Chart {i+1} all methods failed: {str(fallback_error)[:100]}")
                        # Minimal fallback
                        content.append(Paragraph(f"<i>Chart: {chart_title}</i>", normal_style))

                content.append(Spacer(1, 0.4*inch))  # ✅ Professional chart spacing

                if (i + 1) % 2 == 0 and i < len(charts) - 1:
                    content.append(PageBreak())

            except Exception as e:
                print(f"❌ Error processing chart {i+1}: {str(e)}")
                content.append(Paragraph(f"[Chart {i+1}: Processing failed - {str(e)[:100]}]", normal_style))

        print(f"\n📊 Charts Export Summary: {charts_exported}/{total_charts} charts successfully exported")
        if charts_exported < total_charts:
            print(f"⚠️  Warning: {total_charts - charts_exported} chart(s) failed to export")
        elif charts_exported == total_charts and total_charts > 0:
            print(f"✅ Success: All {total_charts} charts exported successfully!")

        # ✅ FIX #19: Quality Score Methodology with icon
        content.append(PageBreak())
        appendix_icon = "📋"
        appendix_title = f"{appendix_icon} Appendix: Quality Score Methodology" if lang == "en" else f"{appendix_icon} Phụ lục: Phương pháp tính Quality Score"
        content.append(Paragraph(appendix_title, heading_style))
        content.append(Spacer(1, 0.2*inch))  # ✅ FIX #17: Consistent spacing

        if lang == "en":
            methodology_text = """
            <b>Based on ISO 8000 Data Quality Standards</b><br/><br/>
            Our Quality Score (0-100) evaluates data across 6 key dimensions:<br/><br/>
            1. <b>Data Completeness (20%)</b>: Percentage of non-null values<br/>
            2. <b>Data Consistency (20%)</b>: Format consistency across columns<br/>
            3. <b>Data Accuracy (20%)</b>: Valid ranges and business rules compliance<br/>
            4. <b>Data Timeliness (15%)</b>: Recency and update frequency<br/>
            5. <b>Data Uniqueness (15%)</b>: Duplicate detection and handling<br/>
            6. <b>Data Validity (10%)</b>: Schema compliance and type correctness<br/><br/>
            <b>Rating Scale:</b><br/>
            • 90-100: [5 STARS] Excellent - Production Ready<br/>
            • 80-89: [4 STARS] Good - Minor improvements recommended<br/>
            • 70-79: [3 STARS] Acceptable - Some issues to address<br/>
            • 60-69: [2 STARS] Fair - Significant improvements needed<br/>
            • 0-59: [1 STAR] Poor - Major data quality issues
            """
        else:
            methodology_text = """
            <b>Dựa trên Tiêu chuẩn Chất lượng Dữ liệu ISO 8000</b><br/><br/>
            Quality Score của chúng tôi (0-100) đánh giá dữ liệu qua 6 tiêu chí chính:<br/><br/>
            1. <b>Độ đầy đủ dữ liệu (20%)</b>: Phần trăm giá trị không null<br/>
            2. <b>Độ nhất quán dữ liệu (20%)</b>: Tính nhất quán định dạng giữa các cột<br/>
            3. <b>Độ chính xác dữ liệu (20%)</b>: Tuân thủ phạm vi hợp lệ và quy tắc nghiệp vụ<br/>
            4. <b>Tính kịp thời (15%)</b>: Độ mới và tần suất cập nhật<br/>
            5. <b>Tính duy nhất (15%)</b>: Phát hiện và xử lý trùng lặp<br/>
            6. <b>Tính hợp lệ (10%)</b>: Tuân thủ schema và đúng kiểu dữ liệu<br/><br/>
            <b>Thang đánh giá:</b><br/>
            • 90-100: [5 SAO] Xuất sắc - Sẵn sàng triển khai<br/>
            • 80-89: [4 SAO] Tốt - Cần cải thiện nhỏ<br/>
            • 70-79: [3 SAO] Chấp nhận được - Một số vấn đề cần giải quyết<br/>
            • 60-69: [2 SAO] Khá - Cần cải thiện đáng kể<br/>
            • 0-59: [1 SAO] Kém - Vấn đề chất lượng dữ liệu nghiêm trọng
            """

        content.append(Paragraph(methodology_text, normal_style))
        content.append(Spacer(1, 0.3*inch))

        # ✅ FIX #13: Limitations and Disclaimers section (transparency & legal protection)
        if lang == "en":
            limitations_title = "Limitations and Disclaimers"
            limitations_text = """
            <b>Data Assumptions and Limitations:</b><br/><br/>
            
            1. <b>Data Quality Dependence</b>: Analysis quality directly depends on input data accuracy. 
            Ensure source data is validated before using insights for critical decisions.<br/><br/>
            
            2. <b>Benchmark Approximations</b>: Industry benchmarks are estimates based on public research 
            (McKinsey, Deloitte, WordStream, etc.) and may not reflect your specific market conditions. 
            Local market variations should be considered.<br/><br/>
            
            3. <b>AI-Generated Insights</b>: Recommendations are AI-assisted and should be reviewed by 
            domain experts before implementation. Human judgment remains essential for strategic decisions.<br/><br/>
            
            4. <b>Currency Assumptions</b>: Currency detection uses heuristic algorithms. Verify currency 
            matches your data (VND vs USD) in the report header. Exchange rates are approximate.<br/><br/>
            
            5. <b>Temporal Scope</b>: Analysis reflects data from the specified date range only. 
            Business conditions may have changed since data collection. Check 'Data Period' in metadata.<br/><br/>
            
            6. <b>Statistical Limitations</b>: Small sample sizes (< 30 rows) may produce unreliable 
            insights. Quality score incorporates sample size considerations but use caution with limited data.<br/><br/>
            
            <b>Disclaimer:</b> This report is for informational purposes only. DataAnalytics Vietnam 
            is not liable for business decisions made based on this analysis. Always validate findings 
            with your team and industry experts before major investments or strategic changes.
            """
        else:
            limitations_title = "Giới Hạn và Miễn Trừ Trách Nhiệm"
            limitations_text = """
            <b>Các Giả Định và Giới Hạn về Dữ Liệu:</b><br/><br/>
            
            1. <b>Phụ thuộc Chất lượng Dữ liệu</b>: Chất lượng phân tích phụ thuộc trực tiếp vào độ 
            chính xác dữ liệu đầu vào. Đảm bảo dữ liệu nguồn được xác thực trước khi sử dụng insights 
            cho quyết định quan trọng.<br/><br/>
            
            2. <b>Benchmark Ước lượng</b>: Các benchmark ngành là ước tính dựa trên nghiên cứu công khai 
            (McKinsey, Deloitte, WordStream, v.v.) và có thể không phản ánh điều kiện thị trường cụ thể 
            của bạn. Cần xem xét biến động thị trường địa phương.<br/><br/>
            
            3. <b>Insights do AI tạo</b>: Khuyến nghị được hỗ trợ bởi AI và nên được xem xét bởi chuyên 
            gia lĩnh vực trước khi triển khai. Phán đoán của con người vẫn thiết yếu cho quyết định chiến lược.<br/><br/>
            
            4. <b>Giả định Tiền tệ</b>: Phát hiện tiền tệ sử dụng thuật toán heuristic. Xác minh tiền tệ 
            khớp với dữ liệu của bạn (VND vs USD) trong header báo cáo. Tỷ giá là ước lượng.<br/><br/>
            
            5. <b>Phạm vi Thời gian</b>: Phân tích chỉ phản ánh dữ liệu từ khoảng thời gian được chỉ định. 
            Điều kiện kinh doanh có thể đã thay đổi kể từ khi thu thập dữ liệu. Kiểm tra 'Chu kỳ dữ liệu' 
            trong metadata.<br/><br/>
            
            6. <b>Giới hạn Thống kê</b>: Kích thước mẫu nhỏ (< 30 hàng) có thể tạo insights không đáng tin cậy. 
            Quality score đã tính đến kích thước mẫu nhưng hãy thận trọng với dữ liệu hạn chế.<br/><br/>
            
            <b>Miễn trừ:</b> Báo cáo này chỉ mang tính chất thông tin. DataAnalytics Vietnam không chịu 
            trách nhiệm cho các quyết định kinh doanh dựa trên phân tích này. Luôn xác thực các phát hiện 
            với đội ngũ và chuyên gia ngành của bạn trước khi đầu tư lớn hoặc thay đổi chiến lược.
            """
        
        # ✅ FIX #19: Add icon to limitations section
        limit_icon = "⚠️"
        limitations_title_with_icon = f"{limit_icon} {limitations_title}"
        content.append(Paragraph(limitations_title_with_icon, heading_style))
        content.append(Spacer(1, 0.2*inch))  # ✅ FIX #17: Consistent spacing
        content.append(Paragraph(limitations_text, normal_style))
        content.append(Spacer(1, 0.3*inch))

        # ✅ FIX #12: Enhanced footer with version control metadata
        content.append(Spacer(1, 0.5*inch))
        footer_style = ParagraphStyle('Footer', parent=normal_style, fontSize=8, textColor=colors.grey, alignment=TA_CENTER, fontName=base_font)
        
        # Version information (ISO 8000-8: Provenance metadata)
        report_version = "v2.0"  # Major release with 13 quality fixes
        generation_timestamp = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
        
        # Extract author/reviewer from domain info if available
        author = result['domain_info'].get('author', 'DataAnalytics AI System')
        expert_role = result['domain_info'].get('expert_role', 'Business Analyst')
        
        # Build comprehensive footer with version control
        footer_text = f"""
        <b>Report Metadata:</b> Version {report_version} | Generated: {generation_timestamp}<br/>
        <b>Author:</b> {author} | <b>Reviewer:</b> {expert_role}<br/>
        <i>Generated by DataAnalytics Vietnam | www.dataanalytics.vn | ISO 8000 Compliant</i>
        """
        
        content.append(Paragraph(footer_text, footer_style))
        
        # Build PDF
        doc.build(content)
        
        # Get PDF bytes
        pdf_bytes = buffer.getvalue()
        buffer.close()
        
        return pdf_bytes
    
    except ImportError:
        raise ImportError("reportlab library required. Install: pip install reportlab kaleido")
    except Exception as e:
        raise Exception(f"PDF export failed: {str(e)}")


def export_to_powerpoint(result: Dict[str, Any], df: Any, lang: str = "vi") -> bytes:
    """
    Export analysis results to professional PowerPoint presentation
    
    Args:
        result: Pipeline result dictionary
        df: Original dataframe
        lang: Language code ('vi' or 'en')
    
    Returns:
        PowerPoint file as bytes
    """
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN
        from pptx.dml.color import RGBColor
        
        # Create presentation
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
        
        # Slide 1: Title
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
        
        # Add title
        title_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1))
        title_frame = title_box.text_frame
        title_para = title_frame.add_paragraph()
        title_para.text = "📊 DATA ANALYSIS REPORT" if lang == "en" else "📊 BÁO CÁO PHÂN TÍCH DỮ LIỆU"
        title_para.font.size = Pt(44)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(30, 64, 175)  # #1E40AF
        title_para.alignment = PP_ALIGN.CENTER
        
        # Add subtitle
        subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(3.5), Inches(8), Inches(0.5))
        subtitle_frame = subtitle_box.text_frame
        subtitle_para = subtitle_frame.add_paragraph()
        subtitle_para.text = "DataAnalytics Vietnam - Professional Business Intelligence"
        subtitle_para.font.size = Pt(18)
        subtitle_para.font.color.rgb = RGBColor(100, 116, 139)
        subtitle_para.alignment = PP_ALIGN.CENTER
        
        # Add date
        date_box = slide.shapes.add_textbox(Inches(1), Inches(6), Inches(8), Inches(0.5))
        date_frame = date_box.text_frame
        date_para = date_frame.add_paragraph()
        date_para.text = datetime.now().strftime("%B %d, %Y")
        date_para.font.size = Pt(14)
        date_para.font.color.rgb = RGBColor(100, 116, 139)
        date_para.alignment = PP_ALIGN.CENTER
        
        # Slide 2: Executive Summary
        slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content
        title = slide.shapes.title
        title.text = "📋 Executive Summary" if lang == "en" else "📋 Tóm Tắt Điều Hành"
        
        content = slide.placeholders[1]
        tf = content.text_frame
        tf.clear()
        
        summary = result['insights'].get('executive_summary', 'No summary available')
        p = tf.add_paragraph()
        p.text = summary
        p.font.size = Pt(16)
        p.level = 0
        
        # Slide 3: Key KPIs
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        title = slide.shapes.title
        title.text = "📈 Key Performance Indicators" if lang == "en" else "📈 Chỉ Số Hiệu Suất Chính"
        
        kpis = result['dashboard'].get('kpis', {})
        if kpis:
            # Add KPIs as table
            rows = min(6, len(kpis)) + 1  # Header + max 6 KPIs
            cols = 4
            
            left = Inches(1)
            top = Inches(2)
            width = Inches(8)
            height = Inches(4)
            
            table = slide.shapes.add_table(rows, cols, left, top, width, height).table
            
            # Header row
            table.cell(0, 0).text = "KPI"
            table.cell(0, 1).text = "Value" if lang == "en" else "Giá trị"
            table.cell(0, 2).text = "Status" if lang == "en" else "Trạng thái"
            table.cell(0, 3).text = "Benchmark" if lang == "en" else "Chuẩn"
            
            # KPI data
            for i, (kpi_name, kpi_info) in enumerate(list(kpis.items())[:6], 1):
                table.cell(i, 0).text = kpi_name
                table.cell(i, 1).text = f"{kpi_info['value']:.1f}"
                table.cell(i, 2).text = kpi_info.get('status', 'N/A')
                table.cell(i, 3).text = str(kpi_info.get('benchmark', 'N/A'))
        
        # Slides 4+: Charts
        charts = result['dashboard']['charts']
        for i, chart in enumerate(charts[:8]):  # Max 8 charts
            slide = prs.slides.add_slide(prs.slide_layouts[5])  # Title only
            title = slide.shapes.title
            title.text = f"📊 Chart {i+1}: {chart.get('title', 'Analysis')}"
            
            try:
                # Convert chart to image
                img_bytes = chart['figure'].to_image(format="png", width=1000, height=600)
                img_stream = io.BytesIO(img_bytes)
                
                # Add image to slide
                left = Inches(1)
                top = Inches(1.5)
                pic = slide.shapes.add_picture(img_stream, left, top, width=Inches(8))
            except Exception as e:
                # Add error text if image generation fails
                text_box = slide.shapes.add_textbox(Inches(2), Inches(3), Inches(6), Inches(1))
                tf = text_box.text_frame
                p = tf.add_paragraph()
                p.text = f"[Chart image generation failed: {str(e)}]"
                p.font.size = Pt(14)
        
        # Slide: Key Insights
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        title = slide.shapes.title
        title.text = "🎯 Key Insights" if lang == "en" else "🎯 Insights Chính"
        
        content = slide.placeholders[1]
        tf = content.text_frame
        tf.clear()
        
        for insight in result['insights'].get('key_insights', [])[:5]:
            impact_emoji = "🔴" if insight['impact'] == 'high' else "🟡" if insight['impact'] == 'medium' else "🟢"
            p = tf.add_paragraph()
            p.text = f"{impact_emoji} {insight['title']}"
            p.font.size = Pt(14)
            p.font.bold = True
            p.level = 0
            
            p2 = tf.add_paragraph()
            p2.text = insight['description']
            p2.font.size = Pt(12)
            p2.level = 1
        
        # Slide: Recommendations
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        title = slide.shapes.title
        title.text = "🚀 Recommendations" if lang == "en" else "🚀 Khuyến Nghị"
        
        content = slide.placeholders[1]
        tf = content.text_frame
        tf.clear()
        
        for rec in result['insights'].get('recommendations', [])[:5]:
            priority_emoji = "🔴" if rec['priority'] == 'high' else "🟡" if rec['priority'] == 'medium' else "🟢"
            p = tf.add_paragraph()
            p.text = f"{priority_emoji} [{rec['priority'].upper()}] {rec['action']}"
            p.font.size = Pt(14)
            p.font.bold = True
            p.level = 0
            
            p2 = tf.add_paragraph()
            p2.text = f"Expected Impact: {rec['expected_impact']} | Timeline: {rec['timeline']}"
            p2.font.size = Pt(11)
            p2.level = 1
        
        # Final Slide: Thank You
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        thank_you_box = slide.shapes.add_textbox(Inches(2), Inches(3), Inches(6), Inches(1))
        tf = thank_you_box.text_frame
        p = tf.add_paragraph()
        p.text = "Thank You!" if lang == "en" else "Cảm ơn!"
        p.font.size = Pt(48)
        p.font.bold = True
        p.font.color.rgb = RGBColor(30, 64, 175)
        p.alignment = PP_ALIGN.CENTER
        
        footer_box = slide.shapes.add_textbox(Inches(1), Inches(6), Inches(8), Inches(0.5))
        tf = footer_box.text_frame
        p = tf.add_paragraph()
        p.text = "DataAnalytics Vietnam | www.dataanalytics.vn | ISO 8000 Compliant"
        p.font.size = Pt(12)
        p.font.color.rgb = RGBColor(100, 116, 139)
        p.alignment = PP_ALIGN.CENTER
        
        # Save to bytes
        ppt_buffer = io.BytesIO()
        prs.save(ppt_buffer)
        ppt_bytes = ppt_buffer.getvalue()
        ppt_buffer.close()
        
        return ppt_bytes
    
    except ImportError:
        raise ImportError("python-pptx library required. Install: pip install python-pptx kaleido")
    except Exception as e:
        raise Exception(f"PowerPoint export failed: {str(e)}")


def create_shareable_link(result: Dict[str, Any], df: Any) -> str:
    """
    Create shareable link for report (placeholder for future implementation)
    
    Args:
        result: Pipeline result
        df: Original dataframe
    
    Returns:
        Shareable URL
    """
    # TODO: Implement cloud storage and link generation
    # For now, return placeholder
    return "https://dataanalytics.vn/reports/[REPORT_ID]"


def send_report_email(email: str, report_bytes: bytes, format: str = "pdf", lang: str = "vi") -> bool:
    """
    Send report via email (placeholder for future implementation)
    
    Args:
        email: Recipient email
        report_bytes: Report file bytes
        format: 'pdf' or 'pptx'
        lang: Language code
    
    Returns:
        Success boolean
    """
    # TODO: Implement email sending with SMTP
    # For now, return placeholder
    return True
